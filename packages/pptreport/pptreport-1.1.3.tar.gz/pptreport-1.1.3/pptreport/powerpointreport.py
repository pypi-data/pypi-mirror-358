import os
import glob
import pprint
import json
import re
import subprocess
import logging
import sys
import tempfile
import fitz
from natsort import natsorted

# Pptx modules
from pptx import Presentation
from pptx.util import Cm

from pptreport.slide import Slide

###############################################################################
# ---------------------------- Helper functions ----------------------------- #
###############################################################################


def _fill_dict(d1, d2):
    """ Fill the keys of d1 with the values of d2 if they are not already present in d1.

    Returns
    --------
    None
        d1 is updated in place.
    """

    for key, value in d2.items():
        if key not in d1:
            d1[key] = value


def _replace_quotes(string):
    """ Replace single quotes with double quotes in a string (such as from the pprint utility to make a valid json file) """

    in_string = False
    for i, letter in enumerate(string):

        if letter == "\"":
            in_string = not in_string  # reverse in_string flag

        elif letter == "'" and in_string is False:  # do not replace single quotes in strings
            string = string[:i] + "\"" + string[i + 1:]  # replace single quote with double quote

    return string


def _convert_to_bool(value):
    """ Convert a value to a boolean type. """

    error_message = f"Could not convert string '{value}' to a boolean value."

    if isinstance(value, bool):  # value is already bool
        return value

    elif isinstance(value, str):
        if value.lower() in ["true", "t", "y", "yes"]:
            return True
        elif value.lower() in ["false", "f", "n", "no"]:
            return False
        else:
            raise ValueError(error_message)
    else:
        raise ValueError(error_message)


def _looks_like_filename(string):
    """ Check if a string looks like a filename.

    Parameters
    ----------
    string : str
        The string to check.

    Returns
    -------
    bool
        True if the string looks like a filename, False otherwise.
    """

    # Check if string is a 1-word string with a dot in it and an extension after the dot
    if len(string.split()) == 1 and "." in string and string.rsplit(".", 1)[1]:
        return True
    else:
        return False  # string does not look like a filename


def _regex_has_group(pattern):
    """ Check if the pattern contains any capturing groups """

    try:
        match = re.search(r'(\(.*?\))', pattern)
        return bool(match)
    except re.error:
        return False


###############################################################################
# -------------------- Class for building presentation ---------------------- #
###############################################################################

class PowerPointReport():
    """ Class for building a PowerPoint presentation """

    _default_slide_parameters = {
        "title": None,
        "slide_layout": 1,
        "content_layout": "grid",
        "content_alignment": "center",
        "outer_margin": 2,
        "inner_margin": 1,
        "left_margin": None,
        "right_margin": None,
        "top_margin": None,
        "bottom_margin": None,
        "n_columns": 2,
        "width_ratios": None,
        "height_ratios": None,
        "notes": None,
        "split": False,
        "show_filename": False,
        "filename_alignment": "center",
        "fill_by": "row",
        "remove_placeholders": False,
        "fontsize": None,
        "pdf_pages": "all",
        "missing_file": "raise",
        "dpi": 300,
        "max_pixels": 1e7,
        "empty_slide": "keep",
        "show_borders": False,
    }

    _valid_slide_parameters = ["content", "grouped_content"] + list(_default_slide_parameters.keys())

    _valid_options = {
        "fill_by": ["row", "column"],
        "missing_file": ["raise", "empty", "text", "skip"],
        "show_filename": ["filename", "filename_ext", "filepath", "filepath_ext", "path"]
    }

    def __init__(self, template=None, size="standard", verbosity=0):
        """ Initialize a presentation object using an existing presentation (template) or from scratch (default) """

        self.template = template
        if template is None:
            self.size = size

        self.global_parameters = None
        self._setup_logger(verbosity)

        self.logger.info("Initializing presentation")
        self._initialize_presentation()

    def _setup_logger(self, verbosity=1):
        """
        Setup a logger for the class.

        Parameters
        ----------
        verbosity : int, default 1
            The verbosity of the logger. 0: ERROR and WARNINGS, 1: INFO, 2: DEBUG

        Returns
        -------
        None
            self.logger is set.
        """

        self.logger = logging.getLogger(self.__class__.__name__)
        self.logger.handlers = []  # remove any existing handlers

        # Test if verbosity is an integer
        try:
            verbosity = int(str(verbosity))  # if verbosity is a bool, converting to str raises an error
        except Exception:
            raise ValueError(f"Verbosity must be an integer - the given value is '{verbosity}'")

        # Setup formatting of handler
        H = logging.StreamHandler(sys.stdout)
        simple_formatter = logging.Formatter("[%(levelname)s] %(message)s")
        debug_formatter = logging.Formatter("[%(levelname)s] [%(name)s:%(funcName)s] %(message)s")

        # Set verbosity and formatting
        if verbosity == 0:
            self.logger.setLevel(logging.WARNING)
            H.setFormatter(simple_formatter)
        elif verbosity == 1:
            self.logger.setLevel(logging.INFO)
            H.setFormatter(simple_formatter)
        elif verbosity == 2:
            self.logger.setLevel(logging.DEBUG)
            H.setFormatter(debug_formatter)
        else:
            raise ValueError("Verbosity must be 0, 1 or 2.")

        self.logger.addHandler(H)

    def _initialize_presentation(self):
        """ Initialize a presentation from scratch. Sets the self._prs and self._slides attributes."""

        self._prs = Presentation(self.template)

        # Get ready to collect configuration
        self._config_dict = {}  # configuration dictionary

        # Set size of the presentation (if not given by a template)
        if self.template is None:
            self.set_size(self.size)  # size is not set if template was given

        # Get ready to add slides
        self._slides = []   # a list of Slide objects

        # Add info to config dict
        if self.template is not None:
            self._config_dict["template"] = self.template
        else:
            self._config_dict["size"] = self.size

    def add_global_parameters(self, parameters):
        """ Add global parameters to the presentation """

        # Test that parameters is a dict
        if not isinstance(parameters, dict):
            raise TypeError(f"Global parameters must be a dictionary. Value given was: {parameters}")

        # Save parameters to self
        self.global_parameters = parameters  # for writing to config file

        # Overwrite default parameters
        for k, v in parameters.items():
            if k not in self._valid_slide_parameters:
                raise ValueError(f"Parameter '{k}' from global parameters is not a valid parameter for slide. Valid parameters are: {self._valid_slide_parameters}")
            else:
                self._default_slide_parameters[k] = v

            if k == "outer_margin":
                self._default_slide_parameters["left_margin"] = v
                self._default_slide_parameters["right_margin"] = v
                self._default_slide_parameters["top_margin"] = v
                self._default_slide_parameters["bottom_margin"] = v

        # Add to internal config dict
        self._config_dict["global_parameters"] = parameters

    def _add_to_config(self, parameters):
        """ Add the slide parameters to the config file. Also checks that the parameters are valid.

        Parameters
        ----------
        parameters : dict
            The parameters for the slide.
        """

        parameters = parameters.copy()  # ensure that later changes in parameters are not reflected in the config dict
        if "slides" not in self._config_dict:
            self._config_dict["slides"] = []

        self._config_dict["slides"].append(parameters)

    def set_size(self, size):
        """
        Set the size of the presentation.

        Parameters
        ----------
        size : str or tuple of float
            Size of the presentation. Can be "standard", "widescreen", "a4-portait" or "a4-landscape". Can also be a tuple of numbers indicating (height, width) in cm.
        """

        if isinstance(size, tuple):
            if len(size) != 2:
                raise ValueError("Size tuple must be of length 2.")
            size = [float(s) for s in size]  # convert eventual strings to floats
            h, w = Cm(size[0]), Cm(size[1])

        elif size == "standard":
            h, w = Cm(19.05), Cm(25.4)

        elif size == "widescreen":
            h, w = Cm(19.05), Cm(33.867)

        elif size == "a4-portrait":
            h, w = Cm(27.517), Cm(19.05)

        elif size == "a4-landscape":
            h, w = Cm(19.05), Cm(27.517)

        else:
            raise ValueError("Invalid size given. Choose from: 'standard', 'widescreen', 'a4-portrait', 'a4-landscape' or a tuple of floats.")

        self._prs.slide_height = h
        self._prs.slide_width = w

    # ------------------------------------------------------ #
    # ------------- Functions for adding slides ------------ #
    # ------------------------------------------------------ #

    def _validate_parameters(self, parameters):
        """ Check the format of the input parameters for the slide and return an updated dictionary. """

        # Set outer margin -> left/right/top/bottom
        orig_parameters = parameters.copy()
        for k in list(parameters.keys()):
            v = orig_parameters[k]

            if k == "outer_margin":
                parameters["left_margin"] = v
                parameters["right_margin"] = v
                parameters["top_margin"] = v
                parameters["bottom_margin"] = v
            else:
                parameters[k] = v  # overwrite previously set top/bottom/left/right margins if they are explicitly given

        # Format integer columns
        int_params = ["n_columns", "max_pixels", "dpi"]
        for param in int_params:
            if param in parameters:
                try:
                    parameters[param] = int(float(str(parameters[param])))  # str() raises an error for bools. float() ensures conversion from 1e10 notation
                except ValueError:
                    raise ValueError(f"Invalid value for '{param}' parameter: '{parameters[param]}'. Please use an integer.")

        # Format "empty_slide"
        if "empty_slide" in parameters:
            if parameters["empty_slide"] not in ["keep", "skip"]:
                raise ValueError(f"Invalid value for 'empty_slide' parameter: {parameters['empty_slide']}. Must be either 'keep' or 'skip'.")

        # Format "split" to int or bool
        if "split" in parameters:
            if not isinstance(parameters["split"], bool):  # only try to convert if not already bool
                try:  # try to convert to int first, e.g. if input is "2"
                    parameters["split"] = int(str(parameters["split"]))
                except Exception:  # if not possible, convert to bool
                    try:
                        parameters["split"] = _convert_to_bool(parameters["split"])
                    except ValueError:
                        raise ValueError(f"Invalid value for 'split' parameter: {parameters['split']}. Must be an integer >= 1 or true/false.")

            if not isinstance(parameters["split"], bool) and parameters["split"] < 1:
                raise ValueError(f"Invalid value for 'split' parameter: {parameters['split']}. Integer must be >= 1.")

        # Format other purely boolean parameters to bool
        bool_parameters = ["remove_placeholders", "show_borders"]
        for param in bool_parameters:
            if param in parameters:
                try:
                    parameters[param] = _convert_to_bool(parameters[param])
                except Exception as e:
                    raise ValueError(f"Invalid value for '{param}' parameter: {parameters[param]}. Error was: {e}")

        # Format show_filename
        if "show_filename" in parameters:
            value = parameters["show_filename"]
            try:
                parameters["show_filename"] = _convert_to_bool(value)
            except Exception:  # if the value is not a bool, it should be a string
                valid = self._valid_options["show_filename"]
                if isinstance(value, str):
                    if value not in valid:
                        raise ValueError(f"Invalid value for 'show_filename' parameter: '{value}'. Please use one of the following: {valid}")
                else:
                    raise ValueError(f"Invalid value for 'show_filename' parameter: '{value}'. Please use one of the following: {valid}")

        # Validate fontsize
        if "fontsize" in parameters:
            try:
                parameters["fontsize"] = float(str(parameters["fontsize"]))  # str ensured that boolean will be an error
            except Exception:
                raise ValueError(f"Invalid value for 'fontsize' parameter: '{parameters['fontsize']}'. Please use a float or integer.")

        # Validate missing_file
        if "missing_file" in parameters:
            parameters["missing_file"] = str(parameters["missing_file"]).lower()
            if parameters["missing_file"] not in self._valid_options["missing_file"]:
                raise ValueError(f"Invalid value for 'missing_file' parameter: '{parameters['missing_file']}'. Must be one of: {self._valid_options['missing_file']}")

        # --- Validate input combinations --- #

        # Establish if content or grouped_content was given
        if "content" in parameters and "grouped_content" in parameters:
            raise ValueError("Invalid input combination. Both 'content' and 'grouped_content' were given - please give only one input type.")

        # If split is given, content should be given
        if parameters.get("split", False) is not False and len(parameters.get("content", [])) == 0:
            raise ValueError("Invalid input combination. 'split' is given, but 'content' is empty")

        # If grouped_content is given, it should be a list
        if "grouped_content" in parameters:
            if not isinstance(parameters["grouped_content"], list):
                raise TypeError(f"Invalid value for 'grouped_content' parameter: {parameters['grouped_content']}. 'grouped_content' must be a list.")

    def add_title_slide(self, title, layout=0, subtitle=None):
        """
        Add a title slide to the presentation.

        Parameters
        ----------
        title : str
            Title of the slide.
        layout : int, default 0
            The layout of the slide. The first layout (0) is usually the default title slide.
        subtitle : str, optional
            Subtitle of the slide if the layout has a suptitle placeholder.
        """

        self.add_slide(title=title, slide_layout=layout)
        slide = self._slides[-1]._slide  # pptx slide object

        # Fill placeholders
        if subtitle is not None:
            if len(slide.placeholders) == 2:
                slide.placeholders[1].text = subtitle

    def add_slide(self,
                  content=None,
                  **kwargs   # arguments given as a dictionary; ensures control over the order of the arguments
                  ):
        """
        Add a slide to the presentation.

        Parameters
        ----------
        content : list of str
            List of content to be added to the slide. Can be either a path to a file or a string.
        grouped_content : list of str
            List of grouped content to be added to the slide. The groups are identified by the regex groups of each element in the list.
        title : str, optional
            Title of the slide.
        slide_layout : int or str, default 1
            Layout of the slide. If an integer, it is the index of the layout. If a string, it is the name of the layout.
        content_layout : str, default "grid"
            Layout of the slide. Can be "grid", "vertical" or "horizontal". Can also be an array of integers indicating the layout of the slide.
        content_alignment : str, default "center"
            Alignment of the content. Can be combinations of "upper", "lower", "left", "right" and "center". Examples: "upper left", "center", "lower right".
            The default is "center", which will align the content centered both vertically and horizontally.
        outer_margin : float, default 2
            Outer margin of the slide (in cm).
        inner_margin : float, default 1
            Inner margin of the slide elements (in cm).
        left_margin / right_margin : float, optional
            Left and right margin of the slide elements (in cm). Can be used to overwrite outer_margin for left/right/both dependent on which are given.
        top_margin / bottom_margin : float, optional
            Top and bottom margin of the slide elements (in cm). Can be used to overwrite outer_margin for top/bottom/both dependent on which are given.
        n_columns : int, default 2
            Number of columns in the layout in case of "grid" layout.
        width_ratios : list of float, optional
            Width of the columns in case of "grid" layout.
        height_ratios : list of float, optional
            Height of the rows in case of "grid" layout.
        notes : str, optional
            Notes for the slide. Can be either a path to a text file or a string.
        split : bool or int, default False
            Split the content into multiple slides. If True, the content will be split into one-slide-per-element. If an integer, the content will be split into slides with that many elements per slide.
        show_filename : bool or str, default False
            Show filenames above images. The style of filename displayed depends on the value given:
            - True or "filename": the filename without path and extension (e.g. "image")
            - "filename_ext": the filename without path but with extension (e.g. "image.png")
            - "filepath": the full path of the image (e.g. "/home/user/image")
            - "filepath_ext": the full path of the image with extension (e.g. "/home/user/image.png")
            - "path": the path of the image without filename (e.g. "/home/user")
            - False: no filename is shown (default)
        filename_alignment : str, default "center"
            Horizontal alignment of the filename. Can be "left", "right" and "center".
            The default is "center", which will align the content centered horizontally.
        fill_by : str, default "row"
            If slide_layout is grid or custom, choose to fill the grid row-by-row or column-by-column. 'fill_by' can be "row" or "column".
        remove_placeholders : str, default False
            Whether to remove empty placeholders from the slide, e.g. if title is not given. Default is False; to keep all placeholders. If True, empty placeholders will be removed.
        fontsize : float, default None
            Fontsize of text content. If None, the fontsize is automatically determined to fit the text in the textbox.
        pdf_pages : int, list of int or "all", default "all"
            Pages to be included from a multipage pdf. e.g. 1 (will include page 1), [1,3] will include pages 1 and 3. "all" includes all available pages.
        missing_file : str, default "raise"
            What to do if no files were found from a content pattern, e.g. "figure*.txt". Can be either "raise", "empty" or "skip".
            - If "raise", a FileNotFoundError will be raised
            - If "text", a content box will be added with the text of the missing content pattern
            - If "empty", an empty content box will be added for the content pattern and 'add_slide' will continue without error
            - If "skip", this content pattern will be skipped (no box added)
        empty_slide : str, "keep" or "skip"
            Whether to keep slides where no file pattern content was found (e.g. if missing_file is "text"/"empty"/"skip" but no files were found for content). Options are:
            - "keep" (default): slides are kept even if the content pattern was not found
            - "skip": slides without any content will not be added to the presentation
        dpi : int, default 300
            Dots per inch of the image. Only used when converting pdf to image.
        max_pixels : int, default 1e7
            Maximum number of pixels in an image. If an image has more pixels than this, it will be resized.
        show_borders : bool, default False
            Whether to show borders of content boxes. This option is useful for debugging layouts. If True, adds a black border to the content boxes. If False, no borders are added.
        """

        self.logger.debug("Started adding slide")

        # Get input parameters;
        parameters = {}
        parameters["content"] = content
        parameters.update(kwargs)
        parameters = {k: v for k, v in parameters.items() if v is not None}
        for param in parameters:
            if param not in self._valid_slide_parameters:
                raise ValueError(f"Invalid parameter '{param}' given for slide. Valid parameters are: {self._valid_slide_parameters}")
        self._add_to_config(parameters)
        self.logger.debug(f"Input parameters: {parameters}")

        # Validate parameters and expand outer_margin
        self._validate_parameters(parameters)  # changes parameters in place

        # If input was None, replace with default parameters from upper presentation
        _fill_dict(parameters, self._default_slide_parameters)
        self.logger.debug("Final slide parameters: {}".format(parameters))

        # Add slides dependent on content type
        if "grouped_content" in parameters:

            content_per_group = self._get_paired_content(parameters["grouped_content"], parameters["missing_file"])

            # If no grouped content was found, add empty slide if empty_slide is "keep"
            if len(content_per_group) == 0 and parameters["empty_slide"] == "keep":
                self.logger.warning(f"No files found for grouped_content: '{parameters['grouped_content']}, but empty_slide == 'keep'. Adding slide without content. Set empty_slide == 'skip' to skip slides without content.")
                slide = self._setup_slide(parameters)

                if parameters["missing_file"] == "text":
                    slide.content = parameters["grouped_content"]
                elif parameters["missing_file"] == "empty":
                    slide.content = [None if _looks_like_filename(element) else element for element in parameters["grouped_content"]]
                else:  # missing_file == "skip"
                    slide.content = [element for element in parameters["grouped_content"] if not _looks_like_filename(element)]  # only skip filenames, not text

                slide._filenames = slide.content
                slide._fill_slide()
                return

            # Create one slide per group
            tmp_files = []
            for group, content in content_per_group.items():

                # Save original filenames / content
                filenames = content[:]

                # Convert pdf to png files
                for idx, element in enumerate(content):
                    if element is not None:  # single files may be missing for groups
                        if element.endswith(".pdf"):
                            img_files = self._convert_pdf(element, parameters["pdf_pages"])

                            if len(img_files) > 1:
                                raise ValueError(f"Invalid value for 'pdf_pages': {parameters['pdf_pages']}. "
                                                 f"Multiple pages in pdf is not supported for grouped content. Found {len(img_files)} pages in {content}, as pdf_pages is set to '{parameters['pdf_pages']}'. "
                                                 "Please adjust pdf_pages to only include one page, e.g. pdf_pages=1.")
                            content[idx] = img_files[0]
                            tmp_files.append(content[idx])

                slide = self._setup_slide(parameters)
                slide.title = f"Group: {group}" if slide.title is None else slide.title
                slide.content = content
                slide._filenames = filenames  # original filenames per content element
                slide._fill_slide()

        else:
            content, filenames, tmp_files = self._get_content(parameters)

            # skip full slide if no content was found and missing_file is set to "skip"
            if content == "skip-slide":
                return  # return before creating slide

            # Create slide(s)
            for i, slide_content in enumerate(content):

                # Setup an empty slide
                slide = self._setup_slide(parameters)
                slide.content = slide_content
                slide._filenames = filenames[i]  # original filenames per content element
                slide._fill_slide()  # Fill slide with content

        # clean tmp files after adding content to slide(s)
        for tmp_file in tmp_files:
            os.remove(tmp_file)

        self.logger.debug("Finished adding slide")

    def _setup_slide(self, parameters):
        """ Initialize an empty slide with a given layout. """

        # How many slides are already in the presentation?
        n_slides = len(self._slides)
        self.logger.info("Adding slide {}".format(n_slides + 1))

        # Add slide to python-pptx presentation
        slide_layout = parameters.get("slide_layout", 0)
        layout_obj = self._get_slide_layout(slide_layout)
        slide_obj = self._prs.slides.add_slide(layout_obj)

        # Add slide to list of slides in internal object
        slide = Slide(slide_obj, parameters)
        slide.logger = self.logger

        # Add information from presentation to slide
        slide._default_parameters = self._default_slide_parameters
        slide._slide_height = self._prs.slide_height
        slide._slide_width = self._prs.slide_width

        self._slides.append(slide)

        return slide

    def _get_slide_layout(self, slide_layout):
        """ Get the slide layout object from a given layout. """

        if isinstance(slide_layout, int):
            try:
                layout_obj = self._prs.slide_layouts[slide_layout]
            except IndexError:
                n_layouts = len(self._prs.slide_layouts)
                raise IndexError(f"Layout index {slide_layout} not found in slide master. The number of slide layouts is {n_layouts} (the maximum index is {n_layouts-1})")

        elif isinstance(slide_layout, str):

            layout_obj = self._prs.slide_layouts.get_by_name(slide_layout)
            if layout_obj is None:
                raise KeyError(f"Layout named '{slide_layout}' not found in slide master.")

        else:
            raise TypeError("Layout should be an integer or a string.")

        return layout_obj

    def _convert_pdf(self, pdf, pdf_pages, dpi=300):
        """ Convert a pdf file to a png file(s).

        Parameters
        ----------
        pdf : str
            pdf file to convert
        pdf_pages: str, int
            pages to include if pdf is a multipage pdf.
            e.g. [1,2] gives firt two pages, all gives all pages
        dpi : int, default 300
            dpi of the output png file

        Returns
        -------
        img_files: [str]
            list containing converted filenames (in the tmp folder)
        """

        # Check that dpi is > 0
        if dpi <= 0:
            raise ValueError(f"Invalid value for 'dpi' parameter: {dpi}. Must be an integer > 0.")

        # open pdf with fitz module from pymupdf
        try:
            def is_pdf(filepath):
                """Check whether file contents follow pdf format."""
                with open(filepath, "rb") as f:
                    header = f.read(5)
                    return header == b"%PDF-"

            if is_pdf(pdf):
                doc = fitz.open(pdf)
            else:
                raise ValueError("File contents not in PDF format!")
        except Exception as e:
            raise ValueError(f"Could not open .pdf file: {pdf}. Error was: {e}")

        # get page count
        pages = doc.page_count
        pages = [i + 1 for i in range(pages)]  # span array over all available pages e.g. pages 3 transforms to [1,2,3]

        # Convert pdf_pages to list
        try:
            pdf_pages = int(str(pdf_pages))  # could be a single value, e.g. 1
        except ValueError:
            try:
                # could be a list of values or strings, e.g. ["1", "2"]
                pdf_pages = [int(str(value)) for value in pdf_pages]
            except ValueError:
                try:
                    if isinstance(pdf_pages, str) and pdf_pages.lower() == "all":  # pdf_pages might be a string "all"
                        pdf_pages = pages
                    elif isinstance(pdf_pages, str) and len(pdf_pages.split(",")) > 1:
                        pdf_pages = [int(str(value)) for value in pdf_pages.split(",")]  # pdf_pages might be a string "1,2" - in this case convert to list of int
                    else:
                        raise ValueError
                except ValueError:
                    raise ValueError(f"Invalid value for 'pdf_pages' parameter: '{pdf_pages}'. Expected an integer, a list of integers or 'all'.")

        if not isinstance(pdf_pages, list):
            pdf_pages = [pdf_pages]

        # all index available? will also fail if index not int
        index_mismatch = [page for page in pdf_pages if page not in pages]
        if len(index_mismatch) != 0:
            raise IndexError(f"Invalid value for 'pdf_pages' parameter: '{pdf_pages}'. Pages {index_mismatch} not available for {pdf}. Available pages are: {pages}.")

        img_files = []
        for page_num in pdf_pages:
            # Create temporary file
            temp_name = next(tempfile._get_candidate_names()) + ".png"
            temp_dir = tempfile.gettempdir()
            temp_file = os.path.join(temp_dir, temp_name)
            self.logger.debug(f"Converting pdf page number{page_num} to temporary png at: {temp_file}")

            # Convert pdf to png
            page = doc.load_page(page_num - 1)  # page 1 is load() with 0
            pix = page.get_pixmap(dpi=dpi)
            pix.save(temp_file)
            img_files.append(temp_file)

        return img_files

    def _expand_files(self, lst, missing_file="raise", empty_slide="keep"):
        """ Expand list of files by unix globbing or regex.

        Parameters
        ----------
        lst : [str]
            list of strings which might (or might not) contain "*" or regex pattern.
        missing_file : str, default "raise"
            What to do if no files are found for a glob pattern.
            - If "raise", a FileNotFoundError will be raised.
            - If "empty", None will be added to the content list.
            - If "skip", this content pattern will be skipped completely.
        empty_slide : str, "keep" or "skip"
            Whether to keep slides where no content was found (e.g. if missing_file is "text"/"empty"/"skip" but no files were found for content).
            Default is "keep". If "skip", empty slides will not be added to the presentation.

        Returns
        -------
        content : [str]
            list of files/content. If empty_slide is "skip", and no content was found for file patterns, the function will return "skip-slide".
        """

        if not isinstance(lst, list):
            lst = [lst]

        patterns = []
        n_patterns = 0        # number of file patterns in lst
        n_patterns_found = 0  # number of file patterns in lst for which files were found

        content = []  # list of files/content
        for element in lst:

            files_found = []   # names of files found for this element

            # If the number of words in element is 1, it could be a file
            if element is not None and len(str(element).split()) == 1:

                element = str(element)  # convert to string, e.g. if element was an int
                element = element.rstrip().lstrip()  # remove trailing and leading spaces to avoid problems with globbing

                # Try to glob files with unix globbing
                globbed = glob.glob(element)
                files_found.extend(globbed)

                # If no files were found by globbing, try to find files by regex
                if len(files_found) == 0:
                    globbed = self._glob_regex(element)
                    files_found.extend(globbed)

                # Establish if the str looks like a filename (or if it is treated a text)
                looks_like_filename = _looks_like_filename(element)
                if looks_like_filename:
                    patterns.append(element)
                    n_patterns += 1

                # Add files to content list if found
                if len(files_found) == 0 and looks_like_filename:  # no files were found, but it looks like a filename
                    if missing_file == "raise":
                        raise FileNotFoundError(f"No files could be found for pattern: '{element}'. Adjust pattern or set missing_file='empty'/'text'/'skip' to ignore the missing file.")
                    elif missing_file == "empty":
                        self.logger.warning(f"No files could be found for pattern: '{element}'. Adding empty box.")
                        content.append(None)
                    elif missing_file == "text":
                        self.logger.warning(f"No files could be found for pattern: '{element}'. Adding textbox with pattern.")
                        content.append(element)
                    elif missing_file == "skip":
                        self.logger.warning(f"No files could be found for pattern: '{element}'. Skipping this file on the slide.")
                    else:
                        raise ValueError(f"Invalid value for 'missing_file' parameter: '{missing_file}'. Must be either 'raise', 'empty', 'text' or 'skip'.")

                elif len(files_found) > 0:
                    n_patterns_found += 1
                    content.append(files_found)

                else:  # no files were found; content is treated as text
                    content.append(element)

            else:  # spaces in text; content is treated as text
                content.append(element)

        # Check if all patterns were found, and whether slide should be skipped
        if n_patterns_found == 0 and n_patterns > 0:
            if empty_slide == "keep":
                self.logger.warning(f"No files were found for any of the filename patterns: {patterns}. Set empty_slide='skip' to skip slides with no content.")
            elif empty_slide == "skip":
                self.logger.info(f"No files were found for any the filename patterns: {patterns}, and empty_slide='skip'. Skipping slide.")
                return "skip-slide"

        # Get the sorted list of files / content
        content_sorted = []  # flattened list of files/content
        for element in content:
            if isinstance(element, list):
                sorted_lst = natsorted(element)
                content_sorted.extend(sorted_lst)
            else:
                content_sorted.append(element)

        return content_sorted

    def _get_content(self, parameters):
        """ Get slide content based on input parameters. """

        # Establish content
        content = parameters.get("content", [])

        # Expand content files
        content = self._expand_files(content, missing_file=parameters["missing_file"], empty_slide=parameters["empty_slide"])
        self.logger.debug(f"Expanded content: {content}")

        # Check if content is empty
        if content == "skip-slide":
            return "skip-slide", None, None

        # Replace multipage pdfs if present
        content_converted = []  # don't alter original list
        filenames = []
        tmp_files = []
        for element in content:
            if isinstance(element, str) and element.endswith(".pdf"):  # avoid None or list type and only replace pdfs
                img_files = self._convert_pdf(element, pdf_pages=parameters.get("pdf_pages", "all"), dpi=parameters.get("dpi", 300))

                content_converted += img_files
                filenames += [element] * len(img_files)  # replace filename with pdf name for each image
                tmp_files += img_files

                self.logger.debug(f"Replaced: {element} with {img_files}.")

            else:
                filenames += [element]
                content_converted += [element]

        content = content_converted

        # If split is false, content should be contained in one slide
        if parameters["split"] is False:
            content = [content]
            filenames = [filenames]
        else:
            content = [content[i:i + parameters["split"]] for i in range(0, len(content), parameters["split"])]
            filenames = [filenames[i:i + parameters["split"]] for i in range(0, len(filenames), parameters["split"])]

        return content, filenames, tmp_files

    def _glob_regex(self, pattern):
        """ Find all files in a directory that match a regex.

        Parameters
        ----------
        pattern : str
            Regex pattern to match files against.

        Returns
        -------
        matched_files : list of str
            List of files that match the regex pattern.
        """

        if pattern is None:
            return []

        # Remove ( and ) from regex as they are only used to group regex later
        pattern_clean = re.sub(r'(?<!\\)[\(\)]', '', pattern)
        self.logger.debug(f"Finding files for possible regex pattern: {pattern_clean}")

        # Find highest existing directory (as some directories might be regex)
        directory = os.path.dirname(pattern_clean)
        while not os.path.exists(directory):
            directory = os.path.dirname(directory)
            if directory == "":
                break  # reached root directory

        # Prepare regex for file search
        pattern_escaped = re.sub(r'(?<!\\)/', r'\\/', pattern)  # Automatically escape / in regex (if not already escaped)
        try:
            pattern_compiled = re.compile(pattern_escaped)
        except re.error:
            # Regex is invalid, assume the pattern is not a regex
            self.logger.warning(f"Pattern is not a valid regex: '{pattern}'. Treating the content as text.")
            return [pattern]
        except Exception as e:
            raise e

        # Find all files that match the regex
        search_glob = os.path.join(directory, "**")
        matched_files = []
        for file in glob.iglob(search_glob, recursive=True):
            if pattern_compiled.match(file):
                matched_files.append(file)

        self.logger.debug(f"Found files: {matched_files}")

        return matched_files

    def _get_paired_content(self, raw_content, missing_file="raise"):
        """ Get content per group from a list of regex patterns.

        Parameters
        ----------
        raw_content : list of str
            List of regex patterns. Each pattern should contain one group.
        missing_file : str, optional
            How to deal with missing files. Options are 'raise', 'empty', 'text' or 'skip'. Default is 'raise'.

        Returns
        -------
        content_per_group : dict
            Dictionary with group names as keys and lists of content as values.
        """

        # Search for regex groups
        group_content = {}  # dict of lists of content input
        for i, pattern in enumerate(raw_content):
            group_content[i] = {}

            # Find all files that match the regex
            files = self._glob_regex(pattern)

            # Check if any files were found and raise error if none were found for file-looking patterns
            if len(files) == 0 and _looks_like_filename(pattern) and missing_file == "raise":
                raise FileNotFoundError(f"No files were found for the pattern: '{pattern}'. Missing_file is set to 'raise'. Adjust the missing_file parameter to 'empty', 'text' or 'skip' to avoid this error.")

            # Find all groups within the regex
            warning = 0
            for fil in files:

                m = re.match(pattern, fil)
                if m:  # if there was a match

                    groups = m.groups()
                    if len(groups) > 1:
                        raise ValueError(f"Invalid value for 'grouped_content' parameter. Regex {pattern} contains more than one group.")
                    elif len(groups) == 1:
                        group = groups[0]
                        group_content[i][group] = fil  # Save the file to the group

                    else:  # 0 groups
                        if warning == 0:
                            s = f"Pattern '{pattern}' does not contain a capturing group (e.g. '(\\w+)_plot.pdf'), but "
                            s += "capturing groups are needed to automatically expand content to multiple slides. "
                            if len(files) > 1:
                                s += f"The pattern matches multiple files, but only one file ('{files[0]}') will be shown. "
                            s += "This content will appear on all expanded slides - please adjust the pattern if needed."
                            self.logger.warning(s)
                            warning += 1  # ensure warning is only printed once per pattern

                        raw_content[i] = files[0]  # replace pattern with file

        # Collect all groups found
        all_regex_groups = sum([list(d.keys()) for d in group_content.values()], [])  # flatten list of lists
        all_regex_groups = natsorted(set(all_regex_groups))
        self.logger.debug(f"Found groups: {all_regex_groups}")

        # If no groups were found for an element, add strings for each group (.e.g. strings/files repeated for each slide)
        content_per_group = {group: [] for group in all_regex_groups}
        for i in group_content:  # index of raw_content
            raw_input = raw_content[i]

            if len(group_content[i]) == 0:  # no groups found
                for group in all_regex_groups:
                    content_per_group[group].append(raw_input)  # this is the same for each slide

            else:
                # Add content for each group with a file
                for group in group_content[i].keys():
                    content_per_group[group].append(group_content[i][group])

                missing_groups = list(set(all_regex_groups) - set(group_content[i].keys()))
                if len(missing_groups) > 0:

                    if missing_file == "raise":
                        s = f"Missing file(s) for grouped content pattern '{raw_input}' for group(s): {missing_groups}."
                        s += " Please ensure that the file(s) exists or adjust 'missing_file' parameter to 'text'/'empty'/'skip' to prevent this error."
                        raise FileNotFoundError(s)

                    else:
                        self.logger.debug(f"Missing file for index {i} for groups: {missing_groups}")
                        for group in missing_groups:
                            pattern = r"\((.*?)\)"
                            element = re.sub(pattern, group, raw_input)
                            content_per_group[group].append(element)  # fill group name into pattern group

        # Check if files are missing for any group
        warnings = []
        for group in content_per_group:
            to_skip = []
            for i, element in enumerate(content_per_group[group]):
                if _looks_like_filename(element) and not os.path.exists(element):
                    if missing_file == "raise":
                        raise FileNotFoundError(f"No file could be found for grouped input: '{element}'. Adjust pattern or set missing_file='empty'/'text'/'skip' to ignore the missing file.")
                    elif missing_file == "empty":
                        s = f"No file could be found for grouped input: '{element}'. Adding empty box."
                        if s not in warnings:  # only print once
                            self.logger.warning(s)
                            warnings.append(s)
                        content_per_group[group][i] = None  # empty box
                    elif missing_file == "skip":
                        s = f"No file could be found for grouped input: '{element}'. Skipping this element."
                        if s not in warnings:  # only print once
                            self.logger.warning(s)
                            warnings.append(s)
                        to_skip.append(i)
                        continue  # skip this element
                    elif missing_file == "text":
                        s = f"No file could be found for grouped input: '{element}'. Adding this element as text."
                        if s not in warnings:  # only print once
                            self.logger.warning(s)
                            warnings.append(s)
                        # keep raw_input as is, as it will be added as text

            for i in to_skip[::-1]:  # reverse order to not mess up indexing
                del content_per_group[group][i]

        # Convert from group per element to element per group
        self.logger.debug(f"Content per group: {content_per_group}")

        return content_per_group

    # ------------------------------------------------------------------------ #
    # --------------------- Saving / loading presentations ---------------------
    # ------------------------------------------------------------------------ #

    def get_config(self, full=False, expand=False):
        """
        Collect a dictionary with the configuration of the presentation

        Parameters
        ----------
        full : bool, default False
            If True, return the full configuration of the presentation. If False, only return the non-default values.
        expand : bool, default False
            If True, expand the content of each slide to a list of files. If False, keep the content as input including "*" and regex.

        Returns
        -------
        config : dict
            Dictionary with the configuration of the presentation.
        """

        # Get configuration of presentation
        if expand is True:  # Read parameters directly from report object

            config = dict(self.__dict__)

            # Add configuration of each slide
            config["slides"] = []
            for slide in self._slides:
                config["slides"].append(slide.get_config())

            # Remove internal variables
            for key in list(config.keys()):  # list to prevent RuntimeError: dictionary changed size during iteration
                if key.startswith("_") or key == "logger":
                    del config[key]

        else:  # Read parameters from internal config_dict
            config = self._config_dict.copy()

        # Get default slide parameters
        defaults = self._default_slide_parameters

        # Resolve configuration of each slide
        for slide_config in config.get("slides", []):
            for key in list(slide_config.keys()):  # list to prevent RuntimeError: dictionary changed size during iteration
                value = slide_config[key]

                # convert bool to str to make it json-compatible
                if isinstance(value, bool):
                    value_converted = str(value)  # convert bool to str to make it json-compatible
                else:
                    value_converted = value
                slide_config[key] = value_converted

                # Remove default values if full is False
                if full is False:
                    if value == defaults.get(key, None):  # compares to the unconverted value
                        del slide_config[key]
                    elif isinstance(value, list) and len(value) == 0:  # content can be an empty list
                        del slide_config[key]

        return config

    def write_config(self, filename, full=False, expand=False):
        """
        Write the configuration of the presentation to a json-formatted file.

        Parameters
        ----------
        filename : str
            Path to the file to write the configuration to.
        full : bool, default False
            If True, write the full configuration of the presentation. If False, only write the non-default values.
        expand : bool, default False
            If True, expand the content of each slide to a list of files. If False, keep the content as input including "*" and regex.
        """

        config = self.get_config(full=full)

        # Get pretty printed config
        pp = pprint.PrettyPrinter(compact=True, sort_dicts=False, width=120)
        config_json = pp.pformat(config)
        config_json = _replace_quotes(config_json)
        config_json = re.sub(r"\"\n\s+\"", "", config_json)  # strings are not allowed to split over multiple lines
        config_json = re.sub(r": None", ": null", config_json)  # Convert to null as None is not allowed in json
        config_json += "\n"  # end with newline

        with open(filename, "w") as f:
            f.write(config_json)

    def from_config(self, config):
        """
        Fill a presentation from a configuration dictionary.

        Parameters
        ----------
        config : str or dict
            A path to a configuration file or a dictionary containing the configuration (such as from Report.get_config()).
        """

        # Load config from file if necessary
        if isinstance(config, str):
            with open(config, "r") as f:
                try:
                    config = json.load(f)
                except Exception as e:
                    raise ValueError("Could not load config file from {}. The error was: {}".format(config, e))

        # Set upper presentation attributes
        upper_keys = config.keys()
        for key in upper_keys:
            if key != "slides":
                setattr(self, key, config[key])

        # Initialize presentation
        self._initialize_presentation()

        # Set global slide parameters
        if "global_parameters" in config:
            self.add_global_parameters(config["global_parameters"])

        # Fill in slides with information from slide config
        for slide_dict in config["slides"]:
            self.add_slide(**slide_dict)  # add all options from slide config
            self.logger.debug("-" * 60)  # separator between slide logging

    def save(self, filename, pdf=False):
        """
        Save the presentation to a file.

        Parameters
        ----------
        filename : str
            Filename of the presentation, e.g. "my_presentation.pptx".
        pdf : bool, default False
            Additionally save the presentation as a pdf file with the same basename as <filename>.
        """

        self.logger.info("Saving presentation to '" + filename + "'")

        # Warning if filename does nto end with .pptx
        if not filename.endswith(".pptx"):
            self.logger.warning("Filename does not end with '.pptx'. This might cause problems when opening the presentation.")

        self._prs.save(filename)

        # Save presentation as pdf
        if pdf:
            self._save_pdf(filename)

    # not included in tests due to libreoffice dependency
    def _save_pdf(self, filename):  # pragma: no cover
        """
        Save presentation as pdf.

        Parameters
        ----------
        filename : str
            Filename of the presentation in pptx format. The pdf will be saved with the same basename.
        """
        self.logger.info("Additionally saving presentation as .pdf")

        # Check if libreoffice is installed
        is_installed = False
        try:
            self.logger.debug("Checking if libreoffice is installed...")
            result = subprocess.run(["libreoffice", "--version"], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            self.logger.debug("Version of libreoffice: " + result.stdout.rstrip())
            is_installed = True

        except FileNotFoundError:
            self.logger.error("Option 'pdf' is set to True, but LibreOffice could not be found on path. Please install LibreOffice to save presentations as pdf.")

        # Save presentation as pdf
        if is_installed:

            outdir = os.path.dirname(filename)
            outdir = "." if outdir == "" else outdir  # outdir cannot be empty

            cmd = f"libreoffice --headless --invisible --convert-to pdf --outdir {outdir} {filename}"
            self.logger.debug("Running command: " + cmd)

            process = subprocess.Popen(cmd, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            while process.poll() is None:
                line = process.stdout.readline().rstrip()
                if line != "":
                    self.logger.debug("Command output: " + line)

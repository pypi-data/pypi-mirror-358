import numpy as np
import os
from pptx.util import Emu, Cm
from pptreport.box import Box
import warnings

# try import for nump>=2 else try numpy<2
try:
    from numpy.exceptions import VisibleDeprecationWarning
except ModuleNotFoundError:
    from numpy import VisibleDeprecationWarning


class Slide():
    """ An internal class for creating slides. """

    def __init__(self, slide, parameters={}):

        self._slide = slide  # Slide object from python-pptx
        self._boxes = {}     # Boxes in the slide; indexed by box id (int)
        self.logger = None

        self.add_parameters(parameters)

    def add_parameters(self, parameters):
        """ Add parameters to the slide as internal variables. """

        for key in parameters:
            if key != "self":
                setattr(self, key, parameters[key])

    def get_config(self):
        """ Get the config dictionary for this slide. """

        config = self.__dict__.copy()  # Make a copy to not change the original dict
        for key in list(config):
            if key.startswith("_") or key == "logger":
                del config[key]

        return config

    def set_layout_matrix(self):
        """ Get the content layout matrix for the slide. """

        # Check validity of n_columns
        try:
            self.n_columns = int(self.n_columns)
            if self.n_columns <= 0:  # Check that n_columns is larger than 0
                raise ValueError
        except ValueError:
            raise ValueError(f"Invalid value for 'n_columns' parameter: '{self.n_columns}'. Please use a positive integer.")

        # Get variables from self
        layout = self.content_layout
        n_elements = len(self.content)
        n_columns = self.n_columns

        # Get layout matrix depending on "layout" variable
        if isinstance(layout, str):
            if layout == "grid":

                n_columns = min(n_columns, n_elements)  # number of columns cannot be larger than number of elements
                n_rows = int(np.ceil(n_elements / n_columns))  # number of rows to fit elements
                n_total = n_rows * n_columns

                intarray = list(range(n_elements))
                intarray.extend([np.nan] * (n_total - n_elements))

                if self.fill_by == "row":
                    layout_matrix = np.array(intarray).reshape((n_rows, n_columns))
                elif self.fill_by == "column":
                    layout_matrix = np.array(intarray).reshape((n_columns, n_rows))
                    layout_matrix = layout_matrix.T
                else:
                    raise ValueError(f"Invalid value for 'fill_by' parameter: '{self.fill_by}'. Please use 'row' or 'column'.")

            elif layout == "vertical":
                layout_matrix = np.array(list(range(n_elements))).reshape((n_elements, 1))

            elif layout == "horizontal":
                layout_matrix = np.array(list(range(n_elements))).reshape((1, n_elements))
            else:
                raise ValueError(f"Invalid value for 'content_layout' parameter: '{layout}'. Please use 'grid', 'vertical' or 'horizontal', or a custom matrix.")

        else:  # layout is expected to be a matrix
            layout_matrix = self._validate_layout(layout)  # check if layout is a valid matrix

        self._layout_matrix = layout_matrix

    # ------------------------ Validate options ------------------------#
    def _validate_margins(self):
        """ Check whether the given margins are valid """

        margins = {"outer_margin": self.outer_margin, "inner_margin": self.inner_margin, "left_margin": self.left_margin, "right_margin": self.right_margin,
                   "top_margin": self.top_margin, "bottom_margin": self.bottom_margin}

        for margin, value in margins.items():
            if value is not None:

                # Check whether value is a float
                try:
                    value = float(value)
                    setattr(self, margin, value)
                except ValueError:
                    raise ValueError(f"Invalid value for '{margin}' parameter: {value}. Could not convert to a float.")

                # Check whether value is positive
                if value < 0:
                    raise ValueError(f"Invalid value for '{margin}' parameter: {value}. Margin cannot be negative")

    def _validate_ratios(self):
        """ Validate the values of width and height ratios """

        parameters = ["width_ratios", "height_ratios"]

        for param in parameters:

            value = getattr(self, param)

            if value is None:
                continue

            # Convert from string to list
            if isinstance(value, str):
                try:
                    value = [float(v) for v in value.split(",")]  # can always be split but not always converted to float
                except Exception:
                    raise ValueError(f"Invalid value for '{param}' parameter: '{value}'. Please give a list of values.")

            # Convert from list of strings to list of floats
            try:
                value = [float(str(v)) for v in value]
                _ = value[0]  # Check that the list is not empty
            except Exception:
                raise ValueError(f"Invalid value for '{param}' parameter: '{value}'. Please give a list of values.")

            # Check that all values are positive
            if any([v <= 0 for v in value]):
                raise ValueError(f"Invalid value for '{param}' parameter: '{value}'. Please give a list of positive values.")

            # Check that the number of values is equal to the number of columns
            expected = self.n_rows if param == "height_ratios" else self.n_cols
            row_col_string = "rows" if param == "height_ratios" else "columns"
            if len(value) > expected:
                self.logger.warning(f"The number of values given in '{param}' is larger than the number of {row_col_string} ({expected}). The extra values will be ignored.")
                value = value[:expected]

            elif len(value) < expected:
                self.logger.warning(f"The number of values given in '{param}' is smaller than the number of {row_col_string} ({expected}). The list will be extended with the last value ({value[-1]}).")
                missing = expected - len(value)
                value.extend([value[-1]] * missing)

            setattr(self, param, value)  # Set the new value

    def _validate_layout(self, layout_matrix):
        """ Validate the given layout matrix. """

        try:
            with warnings.catch_warnings():
                warnings.filterwarnings("error", category=VisibleDeprecationWarning, message="Creating an ndarray from ragged nested*")
                layout_matrix = np.array(layout_matrix)
        except Exception as e:
            raise ValueError(f"Invalid value for 'content_layout' parameter: '{layout_matrix}'. Please make sure that all rows have the same length. Error was: {e}")

        if len(layout_matrix.shape) == 1:
            layout_matrix = layout_matrix.reshape((1, len(layout_matrix)))  # convert to 2D array

        # Make sure that all values are integers
        try:
            layout_matrix = layout_matrix.astype(int)
        except Exception:
            layout_matrix_string = str(layout_matrix).replace('\n', '')  # do not cut numpy array with newlines in error message
            raise ValueError(f"Invalid value for 'content_layout' parameter: '{layout_matrix_string}'. Please use an array of integers.")

        # Check that all values are above -1
        if np.any(layout_matrix < -1):
            layout_matrix_string = str(layout_matrix).replace('\n', '')  # do not cut numpy array with newlines in error message
            raise ValueError(f"Invalid value for 'content_layout' parameter: '{layout_matrix_string}'. Please use an array of integers above -1.")

        # Test that values fit the number of elements
        n_elements = len(self.content)
        content_indices = list(range(n_elements))
        unique_layout_indices = np.unique(layout_matrix)

        # Check that all values in layout_matrix are in content_indices and vice versa
        extra_layout_indices = [i for i in unique_layout_indices if i not in content_indices and i != -1]  # -1 is allowed; will be empty box
        if len(extra_layout_indices) > 0:
            self.logger.warning(f"The layout matrix contains indices that are outside the indices in content (max index = {n_elements-1}): {extra_layout_indices}. These additional content boxes will be empty.")

        missing_layout_indices = [i for i in content_indices if i not in unique_layout_indices]
        if len(missing_layout_indices) > 0:
            content_missing = [self.content[i] for i in missing_layout_indices]
            self.logger.warning(f"Content contains {n_elements} elements, but the content_layout is missing the indices: {missing_layout_indices}. The following contents will not be shown: {content_missing}")

        return layout_matrix

    # -------------------  Fill slide with content  ------------------- #
    def _fill_slide(self):
        """ Fill the slide with content from the internal variables """

        self.set_title()
        self.add_notes()

        # Fill boxes with content
        if len(self.content) > 0:
            self.logger.debug(f"Filling slide with content: {self.content}")
            self.set_layout_matrix()
            self.create_boxes()       # Create boxes based on layout
            self.fill_boxes()         # Fill boxes with content

        # Remove empty placeholders
        if self.remove_placeholders:
            self.remove_empty_ph()

    def set_title(self):
        """ Set the title of the slide. Requires self.title to be set. """

        if self.title is not None:

            # Make sure that title is a string
            self.title = str(self.title)

            if self._slide.shapes.title is None:
                self.logger.warning("Could not set title of slide. The slide does not have a title box.")
            else:
                self._slide.shapes.title.text = self.title

    def add_notes(self):
        """ Add notes to the slide. """

        if self.notes is not None:

            # Convert notes to a list to enable looping
            if not isinstance(self.notes, list):
                self.notes = [self.notes]

            notes_string = ''
            for s in self.notes:
                if isinstance(s, str):
                    if os.path.exists(s):
                        with open(s, 'r') as f:
                            notes_string += f'\n{f.read()}'
                    else:
                        notes_string += f'\n{s}'
                else:
                    raise ValueError("Notes must be either a string or a list of strings.")

            notes_string = notes_string.lstrip()  # remove leading newline
            self._slide.notes_slide.notes_text_frame.text = notes_string

    def create_boxes(self):
        """ Create boxes for the slide dependent on the internal layout matrix. """

        layout_matrix = self._layout_matrix
        nrows, ncols = layout_matrix.shape

        # Check that margins are valid
        self._validate_margins()

        # Establish left/right/top/bottom margins (in cm)
        left_margin = self.outer_margin if self.left_margin is None else self.left_margin
        right_margin = self.outer_margin if self.right_margin is None else self.right_margin
        top_margin = self.outer_margin if self.top_margin is None else self.top_margin
        bottom_margin = self.outer_margin if self.bottom_margin is None else self.bottom_margin
        inner_margin = self.inner_margin

        # Convert margins from cm to pptx units
        left_margin_unit = Cm(left_margin)
        right_margin_unit = Cm(right_margin)
        top_margin_unit = Cm(top_margin)
        bottom_margin_unit = Cm(bottom_margin)
        inner_margin_unit = Cm(inner_margin)

        # Add to top margin based on size of title
        if self._slide.shapes.title.text != "":
            title_margin_unit = Emu(self._slide.shapes.title.top + self._slide.shapes.title.height)
        else:
            title_margin_unit = Emu(0)

        # How many columns and rows are there?
        n_rows, n_cols = layout_matrix.shape
        self.n_rows = n_rows
        self.n_cols = n_cols

        # Get total height and width of pictures
        margin_width = Emu(left_margin_unit + right_margin_unit + (n_cols - 1) * inner_margin_unit)
        margin_height = Emu(top_margin_unit + bottom_margin_unit + (n_rows - 1) * inner_margin_unit + title_margin_unit)
        total_width = self._slide_width - margin_width     # available width for content
        total_height = self._slide_height - margin_height  # available height for content

        # Check if total_width < 0
        if total_width < 0:
            raise ValueError(f"The width of content is negative. The slide width is {self._slide_width.cm:.1f}cm, but the total width of left, right and inner margins is {margin_width.cm}cm. "
                             f"Please adjust the margins to make room for content. Given margins are: "
                             f"left margin={left_margin}cm, right margin={right_margin}cm, inner margin={inner_margin}cm")

        # Check if total_height < 0
        if total_height < 0:
            available_height = self._slide_height.cm - title_margin_unit.cm
            raise ValueError(f"The height of content is negative. The available content height is {available_height:.1f}cm (slide height {self._slide_height.cm:.1f}cm - title height {title_margin_unit.cm:.1f}cm), "
                             f"but the total height of top, bottom and inner margins is {margin_height.cm:.1f}cm. "
                             f"Please adjust the margins to make room for content. Given margins are: "
                             f"top margin={top_margin}cm, bottom margin={bottom_margin}cm, inner margin={inner_margin}cm")

        # Get column widths and row heights
        self._validate_ratios()
        if self.width_ratios is None:
            widths = (np.ones(ncols) / ncols) * total_width
        else:
            widths = np.array(self.width_ratios) / sum(self.width_ratios) * total_width

        if self.height_ratios is None:
            heights = (np.ones(nrows) / nrows) * total_height
        else:
            heights = np.array(self.height_ratios) / sum(self.height_ratios) * total_height

        # Box coordinates
        box_numbers = layout_matrix[~np.isnan(layout_matrix)].flatten().astype(int)
        box_numbers = sorted(set(box_numbers))  # unique box numbers
        box_numbers = [box_number for box_number in box_numbers if box_number >= 0]  # remove negative numbers, e.g. -1 is empty box
        for i in box_numbers:

            # Get column and row number
            coordinates = np.argwhere(layout_matrix == i)

            # Get upper left corner of box
            row, col = coordinates[0]
            left = left_margin_unit + np.sum(widths[:col]) + col * inner_margin_unit
            top = top_margin_unit + np.sum(heights[:row]) + row * inner_margin_unit + title_margin_unit

            # Get total width and height of box (can span multiple columns and rows)
            width = 0
            height = 0
            rows = set(coordinates[:, 0])
            for row in rows:
                height += heights[row]
            height += (len(rows) - 1) * inner_margin_unit  # add inner margins between rows

            cols = set(coordinates[:, 1])
            for col in cols:
                width += widths[col]
            width += (len(cols) - 1) * inner_margin_unit  # add inner margins between columns

            #  Create box
            box = self.get_box((left, top, width, height))

            # Add original filename for the content
            if i < len(self._filenames):  # if i == 2, and number of filenames is 2, index 2 is out of range. Happens if there are empty boxes
                box._filename = self._filenames[i]

            # Add box to dict of boxes
            self._boxes[i] = box

    def get_box(self, coordinates):
        """
        Get a box object with the given coordinates.

        Parameters
        ----------
        coordinates : tuple
            Coordinates containing (left, top, width, height) of the box (in pptx units).
        """

        box = Box(self._slide, coordinates)
        box.logger = self.logger  # share logger with box

        # Add specific parameters to box
        keys = ["content_alignment", "show_filename", "filename_alignment", "fontsize", "max_pixels"]
        parameters = {key: getattr(self, key) for key in keys}
        box.add_parameters(parameters)

        # If show_borders is True for slide
        if self.show_borders:
            box.add_border()

        return box

    def fill_boxes(self):
        """ Fill the boxes with the elements in self.content """

        for i, element in enumerate(self.content):
            if i in self._boxes:  # if there is a box for the element; in custom layouts there can be missing boxes
                self._boxes[i].fill(element, box_index=i)

    def remove_empty_ph(self):
        """ Remove empty placeholders from the slide. """

        if hasattr(self._slide, 'placeholders'):
            for shape in self._slide.placeholders:
                if shape.has_text_frame and shape.text == '':
                    sp = shape.element
                    sp.getparent().remove(sp)
                    self.logger.debug(f"Removed empty placeholder '{shape.name}' (idx: {shape.placeholder_format.idx})")

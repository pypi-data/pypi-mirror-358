import os
import re
import importlib_resources
from contextlib import ExitStack
import tempfile
import numpy as np

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.text.layout import TextFitter
from pptx.oxml.xmlchemy import OxmlElement

from pptreport.config import font_name, md_heading_sizes

# For reading pictures
from PIL import Image
Image.MAX_IMAGE_PIXELS = None  # disable DecompressionBombError


def split_string(string, length):
    """ Split a string into a list of strings of length 'length' """
    return [string[i:i + length] for i in range(0, len(string), length)]


def estimate_fontsize(txt_frame, text, min_size=6, max_size=18, logger=None):
    """
    Resize text to fit the textbox.

    Parameters
    ----------
    txt_frame : pptx.text.text.TextFrame
        The text frame to be resized.
    min_size : int, default 6
        The minimum fontsize of the text.
    max_size : int, default 18
        The maximum fontsize of the text.

    Returns
    --------
    size : int
        The estimated fontsize of the text.
    """

    # Get font
    # https://importlib-resources.readthedocs.io/en/latest/migration.html#pkg-resources-resource-filename
    file_manager = ExitStack()
    ref = importlib_resources.files("pptreport") / "fonts/OpenSans-Regular.ttf"
    font = file_manager.enter_context(importlib_resources.as_file(ref))

    # Calculate best fontsize
    size = None
    try:
        size = TextFitter.best_fit_font_size(text, txt_frame._extents, max_size, font)

    except TypeError:  # happens with long filenames, which cannot fit on one line

        # Try fitting by splitting long words; decrease length if TextFitter still fails
        original_text = text[:]
        max_word_len = 20
        while True:
            if max_word_len < 5:
                if logger is not None:
                    logger.warning(f"Could not fit text '{original_text}' in textbox. Setting fontsize to {min_size}.")
                break  # give up; set text to smallest size
            try:
                words = text.split(" ")
                words = [split_string(word, max_word_len) for word in words]
                words = sum(words, [])  # flatten list
                text = " ".join(words)

                size = TextFitter.best_fit_font_size(text, txt_frame._extents, max_size, font)
                break  # success

            except TypeError:
                max_word_len = int(max_word_len / 2)  # decrease word length
    finally:
        file_manager.close()  # close open font file

    # the output of textfitter is None if the text does not fit; set text to smallest size
    if size is None:
        size = min_size

    # Make sure size is within bounds
    size = max(min_size, size)
    size = min(max_size, size)

    return size


def set_fontsize(run, size):
    """ Set the fontsize of a run.

    Parameters
    ----------
    run : pptx.text.text.Run
        The run to be resized.
    size : int
        The fontsize of the text.
    """

    try:
        run.font.size = Pt(size)
    except Exception as e:
        raise ValueError(f"Invalid value for 'fontsize' parameter: {size}. Error was: {e}")


def set_fontname(run, font_name):
    """ Set the fontname of a run.

    Parameters
    ----------
    run : pptx.text.text.Run
        The run to be resized.
    font_name : str
        The fontname of the text.
    """
    try:
        run.font.name = font_name
    except Exception as e:
        raise ValueError(f"Invalid value for 'fontname' parameter: {font_name}. Error was: {e}")


def set_highlight(run, color):
    """ Set background highlight color of text in run. Method from https://github.com/MartinPacker/md2pptx.

    Color is specified as a hex string, e.g. "#FF0000" for red.
    """

    # get run properties
    rPr = run._r.get_or_add_rPr()

    # Create highlight element
    hl = OxmlElement("a:highlight")

    # Create specify RGB Colour element with color specified
    srgbClr = OxmlElement("a:srgbClr")
    setattr(srgbClr, "val", color)

    # Add colour specification to highlight element
    hl.append(srgbClr)

    # Add highlight element to run properties
    rPr.append(hl)

    return run


def parse_md_structure(data, current_path=[], info={}):
    """ Recursive function to parse the markdown structure from mistune """

    result = []
    if isinstance(data, dict):

        skip_keys = ["type", "children", "text", "alt"]
        info = info.copy()
        info.update({key: data[key] for key in data.keys() if key not in skip_keys})

        if 'type' in data:
            current_path.append(data['type'])
            if 'children' not in data:  # leaf node
                text = data.get('text', data.get('alt', ""))  # text or alt
                result.append((current_path[:], (text, info)))
                current_path.pop()

            else:
                if "children" in data:
                    for child in data["children"]:
                        result.extend(parse_md_structure(child, current_path, info))
                    current_path.pop()

    elif isinstance(data, list):
        for i, item in enumerate(data):
            result.extend(parse_md_structure(item, current_path, info))

    return result


class Box():
    """ A box is a constrained area of the slide which contains a single element e.g. text, a picture, a table, etc. """

    def __init__(self, slide, coordinates):
        """
        Initialize a box.

        Parameters
        ----------
        slide : pptx slide object
            The slide on which the box is located.
        coordinates : tuple
            Coordinates containing (left, top, width, height) of the box (in pptx units).
        """

        self.slide = slide
        self.logger = None

        # Bounds of the box
        self.left = int(coordinates[0])
        self.top = int(coordinates[1])
        self.width = int(coordinates[2])
        self.height = int(coordinates[3])

        # Initialize bounds of the content (can be smaller than the box)
        self.content = None
        self.content_left = self.left
        self.content_top = self.top
        self.content_width = self.width
        self.content_height = self.height

        self.border = None  # border object of the box

    def add_parameters(self, parameters):
        """ Add parameters from the slide """

        for key, value in parameters.items():
            setattr(self, key, value)

    def add_border(self):
        """ Adds a border shape of box to make debugging easier """

        if self.border is None:
            self.border = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.left, self.top, self.width, self.height)
            self.border.shadow.inherit = False  # remove shadow
            self.border.fill.background()
            self.border.line.color.rgb = RGBColor(0, 0, 0)  # black
            self.border.line.width = Pt(1)

    @staticmethod
    def _get_content_type(content):
        """ Determine the type of content. """

        if isinstance(content, str):
            if os.path.isfile(content):
                # Find out the content of the file
                try:
                    with open(content) as f:
                        _ = f.read()
                    return "textfile"

                except UnicodeDecodeError:
                    return "image"
            else:
                return "text"
        elif content is None:
            return "empty"
        else:
            t = type(content)
            return t.__name__

    def fill(self, content, box_index=0):
        """
        Fill the box with content. The function estimates type of content is given.

        Parameters
        ----------
        content : str
            The element to be added to the box.
        box_index : int
            The index of the box (used for accessing properties per box).
        """

        self.box_index = box_index
        self.content = content

        # Find out what type of content it is
        content_type = self._get_content_type(content)
        self.content_type = content_type

        if content_type == "image":
            full_height = self.height

            if self.show_filename is not False:
                # set height of filename to 1/10 of the textbox but at least 290000 (matches Calibri size 12) to ensure the text is still readable
                text_height = max(self.height * 0.1, 290000)
                text_top = self.top
                self.height = full_height - text_height
                self.top = self.top + text_height

            self.fill_image(content)

            # Resize to add filename
            if self.show_filename is not False:
                self.height = text_height  # overwrite height
                self.top = text_top
                _, horizontal = self._get_content_alignment()
                if horizontal != "center":  # make sure text is placed within the picture width, and not the box width, in case alignment is left / right
                    self.left = self.picture.left
                    self.width = self.picture.width

                # Determine filename
                filename = self._filename

                if self.show_filename is True or self.show_filename == "filename":
                    filename = os.path.splitext(os.path.basename(filename))[0]  # basename without extension
                elif self.show_filename == "filename_ext":
                    filename = os.path.basename(filename)  # basename with extension
                elif self.show_filename == "filepath":
                    filename = os.path.splitext(filename)[0]  # filepath without extension
                elif self.show_filename == "filepath_ext":
                    filename = filename  # filepath with extension (original full path)
                elif self.show_filename == "path":
                    filename = os.path.dirname(filename)  # path without filename

                self.fill_text(filename, is_filename=True)

        elif content_type == "textfile":  # textfile can also contain markdown
            with open(content) as f:
                text = f.read()
            self.fill_text(text)

        elif content_type == "text":  # text can also contain markdown
            self.fill_text(content)

        elif content_type == "empty":
            return  # do nothing
        else:
            raise ValueError(f"Content of type '{content_type}' is not supported by pptreport and cannot be added to slide.")

        self.logger.debug(f"Box index {box_index} was filled with {content_type}")

    def fill_image(self, filename):
        """ Fill the box with an image. """

        # Find out the size of the image
        filename, is_temp = self._resize_image(filename)
        self._adjust_image_size(filename)
        self._adjust_image_position()  # adjust image position to middle of box

        # Add image
        self.logger.debug("Adding image to slide from file: " + filename)
        self.picture = self.slide.shapes.add_picture(filename, self.content_left, self.content_top, self.content_width, self.content_height)

        # Remove temporary file
        if is_temp:
            os.remove(filename)

    def _resize_image(self, filename):
        """ Resize the image if needed. Uses max_pixels to determine if image is too large.

        Parameters
        ----------
        filename : str
            Path to the image file.

        Returns
        -------
        tuple : (str, bool)
            The filename of the resized (or original) image and a boolean indicating if the image was resized.
        """

        max_pixels = self.max_pixels
        if max_pixels < 4:
            raise ValueError("Invalid value for 'max_pixels' parameter: '{max_pixels}'. Please use an integer higher than 4.")

        # Check if image can be opened
        try:
            im = Image.open(filename)
        except Exception as e:
            raise ValueError(f"Could not open image file '{filename}'. Error was: {e}")

        im_width, im_height = im.size
        image_pixels = im_width * im_height

        # Resize image if it is too large
        if image_pixels > max_pixels:

            image_ratio = im_width / im_height
            new_height = int(np.sqrt(max_pixels / image_ratio))   # height * height * (width / height) = max_pixels
            new_width = int(new_height * image_ratio)
            self.logger.warning(f"Image '{filename}' is larger than max_pixels={int(max_pixels)} ({im_height}*{im_width}={image_pixels}). Adjust 'max_pixels' to skip resizing. Resizing to size {new_height}*{new_width}...")

            im = im.resize((new_width, new_height), Image.LANCZOS)

            # Create temporary file
            temp_name = next(tempfile._get_candidate_names()) + ".png"
            temp_dir = tempfile.gettempdir()
            temp_file = os.path.join(temp_dir, temp_name)
            im.save(temp_file, quality=100, subsampling=0)

            return temp_file, True  # image was resized

        else:
            return filename, False

    def _adjust_image_size(self, filename):
        """
        Adjust the size of the image to fit the box.

        Parameters
        ----------
        filename : str
            Path to the image file.
        """

        # Find out the size of the image
        im = Image.open(filename)
        im_width, im_height = im.size

        box_width = self.width
        box_height = self.height

        im_ratio = im_width / im_height  # >1 for landscape, <1 for portrait
        box_ratio = box_width / box_height

        # width is the limiting factor; height will be smaller than box_height
        if box_ratio < im_ratio:  # box is wider than image; height will be limiting
            self.content_width = box_width
            self.content_height = box_width * im_height / im_width  # maintain aspect ratio

        # height is the limiting factor; width will be smaller than box_width
        else:
            self.content_width = box_height * im_width / im_height  # maintain aspect ratio
            self.content_height = box_height

    def _get_content_alignment(self):
        """ Get the content alignment for this box. """

        self.logger.debug(f"Getting content alignment for box '{self.box_index}'. Input content alignment is '{self.content_alignment}'")

        # Check if current alignment is valid
        valid_alignments = ["left", "right", "center", "lower", "upper",
                            "lower left", "lower center", "lower right",
                            "upper left", "upper center", "upper right",
                            "center left", "center center", "center right"]

        if isinstance(self.content_alignment, str):  # if content alignment is a string, use it for all boxes
            this_alignment = self.content_alignment

        elif isinstance(self.content_alignment, list):  # if content alignment is a list, use the alignment for the current box
            if self.box_index > len(self.content_alignment) - 1:  # if box index is out of range, use default alignment
                this_alignment = "center"  # default alignment
            else:
                this_alignment = self.content_alignment[self.box_index]
        else:
            raise ValueError(f"Invalid value for 'content_alignment' parameter: '{self.content_alignment}'. Valid content alignments are: {valid_alignments}")

        if this_alignment.lower() not in valid_alignments:
            raise ValueError(f"Invalid value for 'content_alignment' parameter: '{self.content_alignment}'. Valid content alignments are: {valid_alignments}")

        # Expand into the structure "<vertical> <horizontal>"
        if this_alignment.lower() in ["left", "right", "center"]:
            this_alignment = "center " + this_alignment
        elif this_alignment.lower() in ["lower", "upper"]:
            this_alignment = this_alignment + " center"

        return this_alignment.split(" ")

    def _get_filename_alignment(self):
        """ Get the content alignment for this box. """

        self.logger.debug(f"Getting filename alignment for box '{self.box_index}'. Input filename alignment is '{self.filename_alignment}'")

        if isinstance(self.filename_alignment, str):  # if content alignment is a string, use it for all boxes
            this_alignment = self.filename_alignment

        elif isinstance(self.filename_alignment, list):  # if content alignment is a list, use the alignment for the current box
            if self.box_index > len(self.filename_alignment) - 1:  # if box index is out of range, use default alignment
                this_alignment = "center"  # default alignment
            else:
                this_alignment = self.filename_alignment[self.box_index]
        else:
            raise ValueError(f"Invalid value for 'filename_alignment' parameter: {self.filename_alignment}. The input type must be string or list of strings.")

        # Check if current alignment is valid
        valid_alignments = ["left", "right", "center"]
        this_alignment = this_alignment.lower().strip()  # remove any trailing spaces
        if this_alignment not in valid_alignments:
            raise ValueError(f"Invalid value for 'filename_alignment' parameter: {self.filename_alignment}. Valid filename alignments are: {valid_alignments}.")

        return this_alignment

    def _adjust_image_position(self):
        """ Adjust the position of the image to be in the middle of the box. """

        # Get content alignment for this box
        vertical, horizontal = self._get_content_alignment()

        # Adjust image position vertically
        if vertical == "upper":
            self.content_top = self.top
        elif vertical == "lower":
            self.content_top = self.top + self.height - self.content_height
        elif vertical == "center":
            self.content_top = self.top + (self.height - self.content_height) / 2

        # Adjust image position horizontally
        if horizontal == "left":
            self.content_left = self.left
        elif horizontal == "right":
            self.content_left = self.left + self.width - self.content_width
        elif horizontal == "center":
            self.content_left = self.left + (self.width - self.content_width) / 2

    def _fill_md(self, p, text, font_size):
        """
        Fills a paragraph p with basic markdown formatted text, like **Bold**, *italic* ,...
        Supported types:
        - Bold     **bold** / __bold__
        - Italic    *ital*  /  _ital_
        - Link     	[title](https://www.example.com)
        - Heading   #H1 / ## H2 / ...
        (- Image    Only partly - if alternative text is given it will be shown, image should be add via add_image())

        If text does not contain markdown, it will be added as plain text.

        Parameters
        ----------
        p : <pptx.text.text._Paragraph>
            paragraph to add to
        text : str
            The text to be added to the box.
        font_size : int
            The font size of the text. Header sizes are controlled by the fontsizes given in config.md_heading_sizes.
        """

        self.logger.debug(f"Adding markdown formatted text to box '{self.box_index}'.")

        # mistune is only needed for md, only import if needed
        import mistune

        # render input as html.ast
        markdown = mistune.create_markdown(renderer="ast")

        supported_types = ["paragraph", "text", "heading", "newline", "image", "list", "list_item",
                           "block_text", "block_quote", "strong", "emphasis", "link", "codespan", "thematic_break"]

        # traverse the tree
        for i, string in enumerate(text.split("\n")):
            for par in markdown(string):

                # Add newlines between string elements (except for the first element)
                if i > 0:
                    run = p.add_run()
                    run.text = "\n"
                    set_fontsize(run, font_size)

                # Parse all types within this string
                parse_result = parse_md_structure(par, [])

                # Get prefix for the line in case of list / code / quote
                first_element_types = parse_result[0][0]
                if "list" in first_element_types:

                    self.logger.warning("Markdown lists are not supported. Adding element as plain text with number/bullet prefix.")

                    # Get number of whitespaces before the first text
                    prefix = re.match(r"^(\s*).+", string).group(1)  # get number of whitespaces before the first text

                    # Get list item number / bullet point
                    info_dict = parse_result[0][1][1]
                    if info_dict.get("ordered", False):  # ordered list
                        prefix += f"{info_dict.get('start', 1)}. "
                    else:  # unordered list
                        prefix += "• "  # bullet point list

                    parse_result.insert(0, (["paragraph"], (prefix, {})))  # insert prefix as first element (without any formatting)

                elif "thematic_break" in first_element_types:
                    self.logger.warning(f"Markdown horizontal rules are not supported. Adding literal string: '{string}'")
                    parse_result = [(["paragraph"], (string, {}))]

                elif "block_code" in first_element_types:
                    self.logger.warning(f"Markdown code blocks are not supported. Adding literal string: '{string}'")
                    parse_result = [(["paragraph"], (string, {}))]  # replace code block start/end with literal string

                elif "block_quote" in first_element_types:
                    self.logger.warning(f"Markdown block quotes are not supported. Adding literal string: '{string}'")
                    parse_result = [(["paragraph"], (string, {}))]  # replace with literal string

                # Add text to paragraph
                for types, (text, info) in parse_result:
                    # self.logger.debug(f"Adding markdown text to paragraph. Types: {types}. Text: {repr(text)}. Additional info: {info}.")

                    # Check if type is supported
                    for type in types:
                        if type not in supported_types:
                            self.logger.warning(f"pptreport does not support markdown type '{type}' found in string: '{string}'. Adding text instead: '{text}'")
                        elif type == "image":
                            self.logger.warning(f"Markdown images are not supported by pptreport. Adding alternative text instead: '{text}'")

                    # Add run to paragraph
                    run = p.add_run()
                    run.text = text

                    # Adjust formatting
                    run.font.bold = True if "strong" in types else False
                    run.font.italic = True if "emphasis" in types else False

                    # Add highlight for inline code
                    if "codespan" in types:
                        set_highlight(run, "e0e0e0")  # light grey

                    # Add hyperlink
                    if "link" in types:
                        hlink = run.hyperlink
                        hlink.address = info["link"]

                    # Set font size for heading / normal text
                    if "heading" in types:
                        level = info["level"]
                        try:
                            header_fontsize = md_heading_sizes[level]
                        except KeyError:
                            self.logger.warning(f"Header level {level} is not supported. Using default font size.")
                            header_fontsize = font_size
                        set_fontsize(run, header_fontsize)

                    else:
                        set_fontsize(run, font_size)  # font size for plain text

    def fill_text(self, text, is_filename=False):
        """
        Fill the box with text.

        Parameters
        ----------
        text : str
            The text to be added to the box.
        is_filename: bool, optional
            True if text contains a filename to be placed above image, False otherwise. Default: False
        """

        txt_box = self.slide.shapes.add_textbox(self.left, self.top, self.width, self.height)
        txt_frame = txt_box.text_frame
        txt_frame.word_wrap = True

        # Try to fit text size to the box
        if self.fontsize is None:
            self.logger.debug("Estimating fontsize based on text...")
            size = estimate_fontsize(txt_frame, text, logger=self.logger)
            self.logger.debug(f"Fontsize set at: {size}")
        else:
            size = self.fontsize

        # Place all text in one paragraph
        p = txt_frame.paragraphs[0]  # empty paragraph
        self._fill_md(p=p, text=text, font_size=size)  # if text does not contain markdown, it will be added as plain text
        for run in p.runs:
            set_fontname(run, font_name)  # only format font name (size is set during markdown filling)

        # Set alignment of text in textbox
        if is_filename:
            vertical = "lower"
            horizontal = self._get_filename_alignment()
        else:
            vertical, horizontal = self._get_content_alignment()

        if vertical == "upper":
            txt_frame.vertical_anchor = MSO_ANCHOR.TOP
        elif vertical == "lower":
            txt_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        elif vertical == "center":
            txt_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        if horizontal == "left":
            p.alignment = PP_ALIGN.LEFT
        elif horizontal == "right":
            p.alignment = PP_ALIGN.RIGHT
        elif horizontal == "center":
            p.alignment = PP_ALIGN.CENTER

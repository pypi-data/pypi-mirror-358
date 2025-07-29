import csv
import logging
import os
from typing import Dict, Callable, Optional
from pathlib import Path
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from dbfread import DBF
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from pypdf.generic import IndirectObject
import openpyxl
import xlrd
from odf import text, teletype
from odf.opendocument import load
from odf.text import P
import docx
from goblintools.ocr_parser import OCRProcessor

# Document processing libraries (conditional imports)
try:
    import textract

except ImportError as e:
    logging.warning(f"Optional dependency not found: {e}. Some file formats may not be supported.")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TextExtractor:
    """Main class for handling text extraction from various file formats."""
    
    def __init__(self, ocr_handler = False, use_aws=False, aws_access_key=None, aws_secret_key=None, aws_region='us-east-1'):
        """
        Initialize the text extractor.
        
        Args:
            ocr_handler: Optional function to handle OCR for image-based PDFs
        """
        if ocr_handler:
            self.ocr_handler = OCRProcessor(use_aws, aws_access_key, aws_secret_key, aws_region)
        else:
            self.ocr_handler = None

        self._parsers = self._initialize_parsers()

    def _initialize_parsers(self) -> Dict[str, Callable]:
        """Initialize all available text extraction parsers."""
        return {
            '.pdf': self._extract_pdf,
            '.doc': self._extract_doc,
            '.docx': self._extract_docx,
            '.txt': self._extract_txt,
            '.pptx': self._extract_pptx,
            '.html': self._extract_html,
            '.odt': self._extract_odt,
            '.rtf': self._extract_rtf,
            '.csv': self._extract_csv,
            '.xml': self._extract_xml,
            '.xlsx': self._extract_xlsx,
            '.xlsm': self._extract_xlsx,
            '.xls': self._extract_xls,
            '.ods': self._extract_ods,
            '.dbf': self._extract_dbf,
        }

    def add_parser(self, extension: str, parser_func: Callable) -> None:
        """Add or override a parser for a specific file extension."""
        self._parsers[extension.lower()] = parser_func

    def extract_from_file(self, file_path: str) -> str:
        """
        Extract text from a single file.
        
        Args:
            file_path: Path to the file to extract text from
            
        Returns:
            Extracted text as a string
        """
        if not os.path.exists(file_path):
            logger.warning(f"File not found: {file_path}")
            return ""

        file_extension = Path(file_path).suffix.lower()
        parser = self._parsers.get(file_extension)
        
        if not parser:
            logger.warning(f"No parser available for file extension: {file_extension}")
            return ""

        try:
            return parser(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""

    def extract_from_folder(self, folder_path: str) -> str:
        """
        Extract text from all supported files in a folder (recursively).
        
        Args:
            folder_path: Path to the folder to process
            
        Returns:
            Combined extracted text from all files
        """
        if not os.path.exists(folder_path):
            logger.warning(f"Folder not found: {folder_path}")
            return ""

        extracted_text = []
        for root, _, files in os.walk(folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                text = self.extract_from_file(file_path)
                if text:
                    extracted_text.append(text)

        return ' '.join(extracted_text)
    
    def pdf_needs_ocr(self, file_path: str) -> bool:
        with open(file_path, 'rb') as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text and not text.isspace():
                    return False
        return True

    def _resave_pdf(self, file_path: str) -> str:
        from pypdf import PdfReader, PdfWriter

        reader = PdfReader(file_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        output_path = Path(file_path).with_suffix(".resaved.pdf")
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return str(output_path)

    # Individual parser methods
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF files using PyPDF, with fallback to OCR if needed."""
        extracted_text = ''
        has_images = False

        file_path = self._resave_pdf(file_path)

        try:
            reader = PdfReader(file_path)

            for i, page in enumerate(reader.pages):
                try:
                    extracted_text += page.extract_text() or ''

                    resources = page.get('/Resources')
                    if isinstance(resources, IndirectObject):
                        resources = resources.get_object()

                    if not has_images and resources and '/XObject' in resources:
                        xObject = resources['/XObject']
                        if isinstance(xObject, IndirectObject):
                            xObject = xObject.get_object()
                        has_images = any(
                            xObject[obj].get('/Subtype') == '/Image'
                            for obj in xObject
                        )
                except Exception as e:
                    logger.warning(f"Error reading page {i} of {file_path}: {e}")

        except Exception as e:
            logger.error(f"Failed to open PDF {file_path}: {e}")
            return ''

        # Fallback para OCR
        if not extracted_text.strip() and has_images:
            if self.ocr_handler:
                logger.info(f"OCR required for file: {file_path}")
                return self.ocr_handler.extract_text_from_pdf(file_path)
            else:
                logger.warning(f"The file {file_path} requires OCR but no handler was provided.")

        return extracted_text

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        try:
            doc = docx.Document(file_path)
            return ' '.join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {e}")
            return ""

    def _extract_doc(self, file_path: str) -> str:
        """Extract text from legacy DOC files."""
        try:
            return textract.process(file_path).decode('utf-8')
        except Exception as e:
            logger.error(f"Error processing DOC file {file_path}: {e}")
            return ""

    def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error processing TXT file {file_path}: {e}")
                return ""

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint files."""
        try:
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")
            return ""

    def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML files."""
        encodings = ['utf-8', 'latin-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    soup = BeautifulSoup(file.read(), 'html.parser')
                    return soup.get_text(separator=' ', strip=True)
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error processing HTML file {file_path}: {e}")
                break
        return ""

    def _extract_odt(self, file_path: str) -> str:
        """Extract text from OpenDocument Text files."""
        try:
            doc = load(file_path)
            return ' '.join(
                teletype.extractText(element)
                for element in doc.getElementsByType(text.P)
            )
        except Exception as e:
            logger.error(f"Error processing ODT file {file_path}: {e}")
            return ""

    def _extract_rtf(self, file_path: str) -> str:
        """Extract text from RTF files."""
        try:
            with open(file_path, 'r') as file:
                return rtf_to_text(file.read(), errors='ignore')
        except Exception as e:
            logger.error(f"Error processing RTF file {file_path}: {e}")
            return ""

    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return ' '.join(
                    ' '.join(row) 
                    for row in csv.reader(file)
                    if any(field.strip() for field in row)
                )
        except Exception as e:
            logger.error(f"Error processing CSV file {file_path}: {e}")
            return ""

    def _extract_xml(self, file_path: str) -> str:
        """Extract text from XML files."""
        try:
            tree = ET.parse(file_path)
            return ' '.join(
                elem.text.strip() 
                for elem in tree.iter() 
                if elem.text and elem.text.strip()
            )
        except Exception as e:
            logger.error(f"Error processing XML file {file_path}: {e}")
            return ""

    def _extract_xlsx(self, file_path: str) -> str:
        """Extract evaluated text content from Excel files."""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            texts.append(str(cell))
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing XLSX file {file_path}: {e}")
            return ""

    def _extract_xls(self, file_path: str) -> str:
        """Extract text from legacy Excel files."""
        try:
            book = xlrd.open_workbook(file_path, formatting_info=False)
            texts = []
            for sheet in book.sheets():
                for row_idx in range(sheet.nrows):
                    for cell in sheet.row(row_idx):
                        value = cell.value
                        if value and not str(value).startswith('='):
                            texts.append(str(value))
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing XLS file {file_path}: {e}")
            return ""

    def _extract_ods(self, file_path: str) -> str:
        """Extract text from OpenDocument Spreadsheets."""
        try:
            doc = load(file_path)
            return '\n'.join(
                "".join(
                    child.data 
                    for child in p.childNodes 
                    if child.nodeType == child.TEXT_NODE
                )
                for p in doc.getElementsByType(P)
            )
        except Exception as e:
            logger.error(f"Error processing ODS file {file_path}: {e}")
            return ""

    def _extract_dbf(self, file_path: str) -> str:
        """Extract text from DBF database files."""
        try:
            dbf = DBF(file_path, load=True)
            return ' '.join(
                f"{key}: {value}" 
                for record in dbf 
                for key, value in record.items()
            )
        except Exception as e:
            logger.error(f"Error processing DBF file {file_path}: {e}")
            return ""

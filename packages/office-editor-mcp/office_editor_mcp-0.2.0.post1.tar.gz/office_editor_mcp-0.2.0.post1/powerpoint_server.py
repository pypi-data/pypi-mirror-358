"""
MCP Server for PowerPoint Operations

This server provides tools to create, edit and manage PowerPoint presentations.
It's implemented using the Model Context Protocol (MCP) Python SDK.
"""

import os
import sys
import io
import base64
import tempfile
import uuid
import subprocess
import oss2
import re
import random

from pathlib import Path

from loguru import logger
from mcp.server.fastmcp import FastMCP
from mcp.server.fastmcp.prompts import base
from typing import Optional, List, Dict, Any, Union, Tuple

from pptx.slide import Slide
from pptx.text.text import TextFrame

# 标记库是否已安装
pptx_installed = True

# 尝试导入python-pptx库，如果没有安装则标记为未安装但不退出
try:
    import pptx
    from pptx import presentation, Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError:
    print("警告: 未检测到python-pptx库，PowerPoint功能将不可用")
    print("请使用以下命令安装: pip install python-pptx")
    pptx_installed = False

# 尝试导入Pillow库，用于图片处理
pillow_installed = True
try:
    from PIL import Image
except ImportError:
    print("警告: 未检测到Pillow库，图片处理功能将受限")
    print("请使用以下命令安装: pip install Pillow")
    pillow_installed = False

EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')
STEP_BY_STEP_PROCESS_MARKER = '>> '
SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)

# 全局存储内存中的ppt
presentations = {}

# 创建一个MCP服务器，保持名称与配置文件一致
mcp = FastMCP("office editor")


@mcp.tool()
def create_powerpoint_presentation_with_titles(filename: str = "new_presentation", titles: List[str] = None, template: Optional[str] = None) -> str:
    """
    创建一个新的PowerPoint演示文稿。
    
    Args:
        filename: 要创建的文件名 (不需要包含.pptx扩展名)
        titles (List[str]): title slide. Example: ["Title", "Subtitle"]
        template (Optional[str]): The path to the template PPTX file.
                Initializes a presentation from a given template file Or PPTX
                file. (default: :obj:`None`)
    
    Returns:
        prs_id: 该演示文稿在内存中的唯一标识符
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法创建PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    # 确保文件名有.pptx扩展名
    if not filename.lower().endswith('.pptx'):
        filename += '.pptx'

    # Use template if provided, otherwise create new presentation
    if template is not None:
        template_path = Path(template).resolve()
        if not template_path.exists():
            logger.warning(
                f"Template file not found: {template_path}, using "
                "default template"
            )
            prs = Presentation()
        else:
            prs = Presentation(str(template_path))
            # Clear all existing slides by removing them from the slide
            # list
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[-1].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[-1]
    else:
        prs = Presentation()

    prs_id = _generate_presentation_id(filename)
    # 确保id唯一
    while prs_id in presentations:
        prs_id = _generate_presentation_id(filename)

    if titles and isinstance(titles, list):
        title_str = titles[0] if len(titles) > 0 and titles[0] else ""
        subtitle_str = titles[1] if len(titles) > 1 and titles[1] else ""
    else:
        title_str = ""
        subtitle_str = ""

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)

    # Set title and subtitle
    if title_slide.shapes.title:
        title_slide.shapes.title.text_frame.clear()
        _format_text(
            title_slide.shapes.title.text_frame.paragraphs[0],
            title_str,
        )

    if len(title_slide.placeholders) > 1:
        subtitle = title_slide.placeholders[1]
        subtitle.text_frame.clear()
        _format_text(
            subtitle.text_frame.paragraphs[0],
            subtitle_str,
        )

    presentations[prs_id] = prs

    return prs_id


@mcp.tool()
def get_presentation_info(prs_id: str) -> str:
    """
        检查PowerPoint演示文稿的基本信息

        Args:
            prs_id: 该演示文稿在内存中的唯一标识符

        Returns:
            基本信息
        """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法保存PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    if prs_id not in presentations:
        return f"错误：演示文稿ID {prs_id} 未找到"

    try:
        prs = presentations[prs_id]

        # 获取幻灯片信息
        slide_count = len(prs.slides)

        # 获取每张幻灯片的基本信息
        slides_info = []
        for i, slide in enumerate(prs.slides):
            slide_title = "无标题"
            # 尝试获取幻灯片标题
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and (shape.name.startswith("Title") or "标题" in shape.name):
                    slide_title = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                    break

            # 计算幻灯片上的形状数量
            shape_count = len(slide.shapes)

            slides_info.append(f"幻灯片 {i + 1}: {slide_title} (包含 {shape_count} 个形状)")

        # 构建演示文稿信息
        presentation_info = (
                f"prs_id: {prs_id}\n"
                f"幻灯片数量: {slide_count}\n\n"
                f"幻灯片概览:\n" + "\n".join(slides_info)
        )

        return presentation_info
    except Exception as e:
        return f"打开PowerPoint演示文稿时出错: {str(e)}"



# @mcp.tool()
# def add_slide(prs_id: str, layout_name: str = "Title and Content", placeholders: List[str] = None) -> str:
#     """
#     向PowerPoint演示文稿添加新幻灯片。
#
#     Args:
#         prs_id: 该演示文稿在内存中的唯一标识符
#         layout_name: 幻灯片版式名称，常见值包括:
#                     "Title Slide" (标题幻灯片)
#                     "Title and Content" (标题和内容)
#                     "Section Header" (节标题)
#                     "Two Content" (两栏内容)
#                     "Comparison" (比较)
#                     "Title Only" (仅标题)
#                     "Blank" (空白)
#                     "Content with Caption" (带说明的内容)
#                     "Picture with Caption" (带说明的图片)
#         placeholders: 要写入的文本列表，按占位符索引依次填入
#
#     Returns:
#         操作结果信息
#     """
#     # 检查是否安装了必要的库
#     if not pptx_installed:
#         return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
#
#     if prs_id not in presentations:
#         return f"错误：演示文稿ID {prs_id} 未找到"
#
#     prs = presentations[prs_id]
#
#     try:
#         # 获取所有可用的幻灯片版式
#         available_layouts = [layout.name for layout in prs.slide_layouts]
#
#         # 验证提供的版式名称是否有效
#         layout_index = None
#         for i, name in enumerate(available_layouts):
#             if name.lower() == layout_name.lower() or layout_name.lower() in name.lower():
#                 layout_index = i
#                 break
#
#         # 如果找不到匹配的版式，则使用默认版式（通常是"Title and Content"或index=1）
#         if layout_index is None:
#             # 尝试查找最接近的版式
#             if "title" in layout_name.lower() and "content" in layout_name.lower():
#                 # 尝试查找标题和内容版式
#                 for i, name in enumerate(available_layouts):
#                     if "title" in name.lower() and "content" in name.lower():
#                         layout_index = i
#                         break
#
#             # 如果仍然找不到，使用索引1（通常是标题和内容）或0（标题）
#             if layout_index is None:
#                 if len(prs.slide_layouts) > 1:
#                     layout_index = 1  # 标题和内容
#                 else:
#                     layout_index = 0  # 默认为第一个布局
#
#             layout_name = available_layouts[layout_index]
#
#         # 添加新幻灯片
#         slide_layout = prs.slide_layouts[layout_index]
#         new_slide = prs.slides.add_slide(slide_layout)
#
#         for i, text in enumerate(placeholders):
#             if i >= len(new_slide.placeholders):
#                 break
#             try:
#                 new_slide.placeholders[i].text = text
#             except Exception as e:
#                 print(f"占位符{i}无法写入{e}")
#
#         # 获取当前幻灯片总数
#         total_slides = len(prs.slides)
#
#         return (f"已成功添加新幻灯片（版式: {layout_name}）\n"
#                 f"并填充了{len(placeholders)}个占位符"
#                 f"当前演示文稿共有 {total_slides} 张幻灯片")
#     except Exception as e:
#         return f"添加幻灯片时出错: {str(e)}"


@mcp.tool()
def add_slides(prs_id: str, content: Union[str, List[Dict]]) -> str:
    """
    批量添加幻灯片并为每页填充占位符内容。

    Args:
        prs_id: PowerPoint 演示文稿的唯一标识
        content (str): The content to write to the PPTX file as a JSON
                string. Must represent a list of dictionaries with the
                following structure:
                - Dicts: content slides, which can be one of:
                    * Bullet/step slides: {"heading": str, "bullet_points":
                    list of str or nested lists, "img_keywords": str
                    (optional)}
                        - If any bullet point starts with '>> ', it will be
                        rendered as a step-by-step process.
                        - "img_keywords" can be a URL or search keywords for
                        an image (optional).
                    * Table slides: {"heading": str, "table": {"headers": list
                    of str, "rows": list of list of str}}

    Returns:
        操作结果字符串

    Example:
        [
            {
                "heading": "Slide Title",
                "bullet_points": [
                    "**Bold text** for emphasis",
                    "*Italic text* for additional emphasis",
                    "Regular text for normal content"
                ],
                "img_keywords": "relevant search terms for images"
            },
            {
                "heading": "Step-by-Step Process",
                "bullet_points": [
                    ">> **Step 1:** First step description",
                    ">> **Step 2:** Second step description",
                    ">> **Step 3:** Third step description"
                ],
                "img_keywords": "process workflow steps"
            },
            {
                "heading": "Comparison Table",
                "table": {
                    "headers": ["Column 1", "Column 2", "Column 3"],
                    "rows": [
                        ["Row 1, Col 1", "Row 1, Col 2", "Row 1, Col 3"],
                        ["Row 2, Col 1", "Row 2, Col 2", "Row 2, Col 3"]
                    ]
                },
                "img_keywords": "comparison visualization"
            }
        ]
    """
    if prs_id not in presentations:
        return f"错误: 未找到演示文稿 {prs_id}"

    prs = presentations[prs_id]

    logger.info(f"content: {content}\n"
                f"content_type: {type(content)}")

    # Parse and validate content format
    if isinstance(content, str):
        try:
            import json

            parsed_content = json.loads(content)
        except json.JSONDecodeError as e:
            logger.error(f"Content must be valid JSON: {e}")
            return "Failed to parse content as JSON"
    elif isinstance(content, list):
        parsed_content = content
    else:
        logger.error(f"Unsupported content type: {type(content).__name__}")
        return "Content must be a JSON string or a list of dictionaries"

    if not isinstance(parsed_content, list):
        logger.info(f"parsed_content: {parsed_content}\n"
                    f"parsed_content: {type(parsed_content)}")
        logger.error(
            f"PPTX content must be a list of dictionaries, "
            f"got {type(parsed_content).__name__}"
        )
        return "PPTX content must be a list of dictionaries"

    slide_width_inch, slide_height_inch = (
        _get_slide_width_height_inches(prs)
    )

    if parsed_content:
        for slide_data in parsed_content:
            if not isinstance(slide_data, dict):
                continue

            # Handle different slide types
            if 'table' in slide_data:
                _handle_table(
                    prs,
                    slide_data,
                )
            elif 'bullet_points' in slide_data:
                if any(
                    step.startswith(STEP_BY_STEP_PROCESS_MARKER)
                    for step in slide_data['bullet_points']
                ):
                    _handle_step_by_step_process(
                        prs,
                        slide_data,
                        slide_width_inch,
                        slide_height_inch,
                    )
                else:
                    _handle_default_display(
                        prs,
                        slide_data,
                    )
    return "Successfully added slides"


def _handle_table(prs: "presentation.Presentation", slide_json: Dict[str, Any]) -> None:
    r"""Add a table to a slide.

    Args:
        prs (presentation.Presentation): The presentation object.
        slide_json (Dict[str, Any]): The content of the slide as JSON data.
    """
    headers = slide_json['table'].get('headers', [])
    rows = slide_json['table'].get('rows', [])
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = _remove_slide_number_from_heading(
        slide_json['heading']
    )
    left = slide.placeholders[1].left
    top = slide.placeholders[1].top
    width = slide.placeholders[1].width
    height = slide.placeholders[1].height
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height
    ).table

    # Set headers
    for col_idx, header_text in enumerate(headers):
        table.cell(0, col_idx).text = header_text
        table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True

    # Fill in rows
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(row_data):
            table.cell(row_idx, col_idx).text = cell_text


def _handle_default_display(prs: "presentation.Presentation", slide_json: Dict[str, Any]) -> None:
    r"""Display a list of text in a slide.

    Args:
        prs (presentation.Presentation): The presentation object.
        slide_json (Dict[str, Any]): The content of the slide as JSON data.
    """
    status = False

    if 'img_keywords' in slide_json:
        status = _handle_display_image__in_foreground(
            prs,
            slide_json,
        )

    if status:
        return

    #No image, display text only
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    shapes = slide.shapes
    title_shape = shapes.title

    try:
        body_shape = shapes.placeholders[1]
    except KeyError:
        # Get placeholders from the slide without layout_number
        placeholders = _get_slide_placeholders(slide)
        body_shape = shapes.placeholders[placeholders[0][0]]

    title_shape.text = _remove_slide_number_from_heading(
        slide_json['heading']
    )
    text_frame = body_shape.text_frame

    flat_items_list = _get_flat_list_of_contents(
        slide_json['bullet_points'], level=0
    )
    _add_bulleted_items(text_frame, flat_items_list)


def _handle_display_image__in_foreground(prs: "presentation.Presentation", slide_json: Dict[str, Any]) -> bool:
    r"""Create a slide with text and image using a picture placeholder
    layout.

    Args:
        prs (presentation.Presentation): The presentation object.
        slide_json (Dict[str, Any]): The content of the slide as JSON data.

    Returns:
        bool: True if the slide has been processed.
    """
    from io import BytesIO
    import requests

    img_keywords = slide_json.get('img_keywords', '').strip()
    slide = prs.slide_layouts[8]  # Picture with Caption
    slide = prs.slides.add_slide(slide)
    placeholders = None

    title_placeholder = slide.shapes.title  # type: ignore[attr-defined]
    title_placeholder.text = _remove_slide_number_from_heading(
        slide_json['heading']
    )

    try:
        pic_col = slide.shapes.placeholders[1]  # type: ignore[attr-defined]
    except KeyError:
        # Get placeholders from the slide without layout_number
        placeholders = self._get_slide_placeholders(slide)  # type: ignore[arg-type]
        pic_col = None
        for idx, name in placeholders:
            if 'picture' in name:
                pic_col = slide.shapes.placeholders[idx]  # type: ignore[attr-defined]

    try:
        text_col = slide.shapes.placeholders[2]  # type: ignore[attr-defined]
    except KeyError:
        text_col = None
        if not placeholders:
            placeholders = self._get_slide_placeholders(slide)  # type: ignore[arg-type]

        for idx, name in placeholders:
            if 'content' in name:
                text_col = slide.shapes.placeholders[idx]  # type: ignore[attr-defined]

    flat_items_list = _get_flat_list_of_contents(
        slide_json['bullet_points'], level=0
    )
    _add_bulleted_items(text_col.text_frame, flat_items_list)

    if not img_keywords:
        return True

    if isinstance(img_keywords, str) and img_keywords.startswith(
            ('http://', 'https://')
    ):
        try:
            img_response = requests.get(img_keywords, timeout=30)
            img_response.raise_for_status()
            image_data = BytesIO(img_response.content)
            pic_col.insert_picture(image_data)
            return True
        except Exception as ex:
            logger.error(
                'Error while downloading image from URL: %s', str(ex)
            )

    try:
        url = 'https://api.pexels.com/v1/search'
        api_key = os.getenv('PEXELS_API_KEY')

        headers = {
            'Authorization': api_key,
            'User-Agent': 'Mozilla/5.0 (X11; Linux x86_64; rv:10.0) '
                          'Gecko/20100101 Firefox/10.0',
        }
        params = {
            'query': img_keywords,
            'size': 'medium',
            'page': 1,
            'per_page': 3,
        }
        response = requests.get(
            url, headers=headers, params=params, timeout=12
        )
        response.raise_for_status()
        json_response = response.json()

        if json_response.get('photos'):
            photo = random.choice(json_response['photos'])
            photo_url = photo.get('src', {}).get('large') or photo.get(
                'src', {}
            ).get('original')

            if photo_url:
                # Download and insert the image
                img_response = requests.get(
                    photo_url, headers=headers, stream=True, timeout=12
                )
                img_response.raise_for_status()
                image_data = BytesIO(img_response.content)

                pic_col.insert_picture(image_data)
    except Exception as ex:
        logger.error(
            'Error occurred while adding image to slide: %s', str(ex)
        )

    return True

def _handle_step_by_step_process(prs: "presentation.Presentation", slide_json: Dict[str, Any], slide_width_inch: float, slide_height_inch: float) -> None:
    r"""Add shapes to display a step-by-step process in the slide.

    Args:
        prs (presentation.Presentation): The presentation object.
        slide_json (Dict[str, Any]): The content of the slide as JSON data.
        slide_width_inch (float): The width of the slide in inches.
        slide_height_inch (float): The height of the slide in inches.
    """
    import pptx
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches, Pt

    steps = slide_json['bullet_points']
    n_steps = len(steps)

    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = _remove_slide_number_from_heading(
        slide_json['heading']
    )

    if 3 <= n_steps <= 4:
        # Horizontal display
        height = Inches(1.5)
        width = Inches(slide_width_inch / n_steps - 0.01)
        top = Inches(slide_height_inch / 2)
        left = Inches(
            (slide_width_inch - width.inches * n_steps) / 2 + 0.05
        )

        for step in steps:
            shape = shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height
            )
            text_frame = shape.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = pptx.enum.text.PP_ALIGN.CENTER
            text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
            _format_text(
                paragraph, step.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
            )
            for run in paragraph.runs:
                run.font.size = Pt(14)
            left = Inches(left.inches + width.inches - Inches(0.4).inches)
    elif 4 < n_steps <= 6:
        # Vertical display
        height = Inches(0.65)
        top = Inches(slide_height_inch / 4)
        left = Inches(1)
        width = Inches(slide_width_inch * 2 / 3)

        for step in steps:
            shape = shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.PENTAGON, left, top, width, height
            )
            text_frame = shape.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = pptx.enum.text.PP_ALIGN.CENTER
            text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
            _format_text(
                paragraph, step.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
            )
            for run in paragraph.runs:
                run.font.size = Pt(14)
            top = Inches(top.inches + height.inches + Inches(0.3).inches)
            left = Inches(left.inches + Inches(0.5).inches)


@mcp.tool()
def remove_slide(prs_id: str, slide_index: int) -> str:
    """
    Remove a slide from the presentation.

    Args:
        prs_id: The presentation ID
        slide_index: The index of the slide to remove (1-based)

    Returns:
        Confirmation message
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"

    prs = presentations[prs_id]

    # Convert 1-based index to 0-based index
    zero_based_index = slide_index - 1

    if zero_based_index < 0 or zero_based_index >= len(prs.slides):
        return f"Error: Slide index {slide_index} is out of range (valid range: 1-{len(prs.slides)})"

    try:
        # This is a simplified approach that might not work in all python-pptx versions
        # Get XML ids needed for removal
        if hasattr(prs.slides, '_sldIdLst'):
            slide_id = prs.slides._sldIdLst[zero_based_index].attrib['id']
            slide_rel = prs.slides._sldIdLst[zero_based_index].attrib[
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']

            # Remove the slide
            prs.part.drop_rel(slide_rel)
            del prs.slides._sldIdLst[zero_based_index]
        else:
            return f"Cannot remove slide - this functionality requires direct XML manipulation that may not be supported in this version of python-pptx"
    except Exception as e:
        print(f"Error removing slide: {e}", file=sys.stderr)
        return f"Error removing slide: {str(e)}"

    return f"Removed slide at position {slide_index}"


@mcp.tool()
def reorder_slides(prs_id: str, slide_index: int, new_position: int) -> str:
    """
    调整PowerPoint演示文稿中幻灯片的顺序（内存操作）。

    Args:
        prs_id: 演示文稿ID
        slide_index: 要移动的幻灯片索引（从1开始计数）
        new_position: 幻灯片的新位置（从1开始计数）

    Returns:
        操作结果信息
    """
    if not pptx_installed:
        return "错误: 请先安装python-pptx库: pip install python-pptx"

    if prs_id not in presentations:
        return f"错误: 演示文稿ID {prs_id} 未找到"

    try:
        prs = presentations[prs_id]
        total_slides = len(prs.slides)

        # 检查索引有效性
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"

        if new_position < 1 or new_position > total_slides:
            return f"错误: 无效的目标位置 {new_position}，演示文稿共有 {total_slides} 张幻灯片"

        if slide_index == new_position:
            return f"幻灯片 {slide_index} 已经在位置 {new_position}，无需调整"

        # 直接操作XML结构调整顺序
        sldIdLst = prs.slides._sldIdLst
        slide_to_move = sldIdLst[slide_index - 1]

        # 移除原位置的幻灯片
        del sldIdLst[slide_index - 1]

        # 插入到新位置
        # 注意：如果新位置在原位置之后，索引需要减1
        insert_index = new_position - 1 if new_position < slide_index else new_position - 1
        sldIdLst.insert(insert_index, slide_to_move)

        return f"已成功将幻灯片 {slide_index} 移动到位置 {new_position}"

    except Exception as e:
        # 如果直接XML操作失败，使用更安全的复制方法
        return safe_reorder_slides(prs_id, slide_index, new_position)


def safe_reorder_slides(prs_id: str, slide_index: int, new_position: int) -> str:
    """
    调整PowerPoint演示文稿中幻灯片的顺序。
    
    Args:
        prs_id: 演示文稿在内存中的唯一标识符
        slide_index: 要移动的幻灯片索引（从1开始计数）
        new_position: 幻灯片的新位置（从1开始计数）
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"

    prs = Presentation(prs_id)
    total_slides = len(prs.slides)

    if slide_index < 1 or slide_index > total_slides:
        return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"

    if new_position < 1 or new_position > total_slides:
        return f"错误: 无效的目标位置 {new_position}，演示文稿共有 {total_slides} 张幻灯片"

    # 如果位置相同，无需调整
    if slide_index == new_position:
        return f"幻灯片 {slide_index} 已经在位置 {new_position}，无需调整"

    try:
        temp_prs = Presentation()

        # python-pptx目前不直接支持重新排序幻灯片
        # 需要创建一个新的演示文稿并按照新的顺序复制幻灯片

        # 调整幻灯片顺序
        slides = list(prs.slides)

        # 从列表中删除要移动的幻灯片
        slide_to_move = slides.pop(slide_index - 1)

        # 在新位置插入幻灯片
        slides.insert(new_position - 1, slide_to_move)

        # 按新顺序复制幻灯片
        for slide in slides:
            # 尝试获取原始幻灯片的布局
            try:
                original_layout_idx = slide.slide_layout.index
                if original_layout_idx < len(temp_prs.slide_layouts):
                    slide_layout = temp_prs.slide_layouts[original_layout_idx]
                else:
                    slide_layout = temp_prs.slide_layouts[0]  # 默认使用第一个布局
            except:
                slide_layout = temp_prs.slide_layouts[0]  # 如果失败，使用第一个布局

            # 创建新幻灯片
            new_slide = temp_prs.slides.add_slide(slide_layout)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    # 复制文本框
                    new_shape = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    new_shape.text_frame.text = shape.text_frame.text
                elif shape.shape_type == MSO_SHAPE.PICTURE:
                    # 复制图片
                    try:
                        image_data = shape.image.blob
                        image_stream = io.BytesIO(image_data)
                        new_slide.shapes.add_picture(
                            image_stream, shape.left, shape.top, shape.width, shape.height
                        )
                    except Exception:
                        # 简单回退方案
                        new_shape = new_slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            shape.left, shape.top, shape.width, shape.height
                        )
                        new_shape.text = "[图片]"
                elif shape.shape_type == MSO_SHAPE.TABLE:
                    # 复制表格
                    table = shape.table
                    new_table = new_slide.shapes.add_table(
                        table.rows.count, table.columns.count,
                        shape.left, shape.top, shape.width, shape.height
                    ).table

                    # 复制表格内容
                    for r in range(table.rows.count):
                        for c in range(table.columns.count):
                            cell = table.cell(r, c)
                            new_cell = new_table.cell(r, c)
                            new_cell.text = cell.text
                else:
                    # 其他形状的简单复制
                    new_shape = new_slide.shapes.add_shape(
                        shape.shape_type,
                        shape.left, shape.top, shape.width, shape.height
                    )
                    if shape.has_text_frame:
                        new_shape.text_frame.text = shape.text_frame.text

        # 保存临时演示文稿
        presentations[prs_id] = temp_prs
        return f"已成功将幻灯片 {slide_index} 移动到位置 {new_position}"

    except Exception as e:
        return f"调整幻灯片顺序时出错: {str(e)}"


@mcp.tool()
def set_slide_background(prs_id: str, slide_index: int, background_color: str = None) -> str:
    """
    为PowerPoint演示文稿中的指定幻灯片设置背景颜色。
    
    Args:
        prs_id: PowerPoint演示文稿在内存中的唯一标识符
        slide_index: 要设置背景的幻灯片索引（从1开始计数），0表示所有幻灯片
        background_color: 背景颜色的十六进制代码，例如"#FF0000"表示红色
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    # 确保文件存在
    if prs_id not in presentations:
        return f"错误: 演示文稿ID {prs_id} 未找到"

    # 验证颜色格式
    if background_color:
        if not background_color.startswith('#') or len(background_color) != 7:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"

        try:
            # 解析RGB值
            r = int(background_color[1:3], 16)
            g = int(background_color[3:5], 16)
            b = int(background_color[5:7], 16)
        except ValueError:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"

    try:
        # 打开PowerPoint演示文稿
        prs = presentations[prs_id]

        # 获取幻灯片总数
        total_slides = len(prs.slides)

        # 如果slide_index不是0，检查其有效性
        if slide_index != 0 and (slide_index < 1 or slide_index > total_slides):
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"

        # 确定要修改的幻灯片
        slides_to_modify = []
        if slide_index == 0:
            # 修改所有幻灯片
            slides_to_modify = list(range(total_slides))
        else:
            # 修改指定幻灯片
            slides_to_modify = [slide_index - 1]  # 转换为0开始的索引

        # 设置背景颜色
        for idx in slides_to_modify:
            slide = prs.slides[idx]

            # 获取幻灯片背景
            background = slide.background

            # 设置背景填充类型为纯色
            fill = background.fill
            fill.solid()

            # 设置背景颜色
            if background_color:
                fill.fore_color.rgb = RGBColor(r, g, b)

        # 构建结果消息
        if slide_index == 0:
            return f"已成功设置所有幻灯片的背景颜色为 {background_color}"
        else:
            return f"已成功设置幻灯片 {slide_index} 的背景颜色为 {background_color}"

    except Exception as e:
        return f"设置幻灯片背景时出错: {str(e)}"


@mcp.tool()
def add_text_box(
        prs_id: str,
        slide_index: int,
        text: str,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 4.0,
        height: float = 1.0,
        font_name: str = None,
        font_size: int = None,
        font_bold: bool = False,
        font_italic: bool = False,
        text_color: str = None,
        alignment: str = "left"
) -> str:
    """
    向PowerPoint演示文稿的指定幻灯片添加文本框。
    
    Args:
        prs_id: PowerPoint演示文稿在内存中的唯一标识
        slide_index: 要添加文本框的幻灯片索引（从1开始计数）
        text: 要添加的文本内容
        left: 文本框左侧位置（英寸）
        top: 文本框顶部位置（英寸）
        width: 文本框宽度（英寸）
        height: 文本框高度（英寸）
        font_name: 字体名称，如"宋体"、"Arial"等
        font_size: 字体大小（磅）
        font_bold: 是否加粗
        font_italic: 是否斜体
        text_color: 文本颜色的十六进制代码，例如"#FF0000"表示红色
        alignment: 文本对齐方式，可选值: "left", "center", "right", "justify"
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    # 确保文件存在
    if prs_id not in presentations:
        return f"错误: 演示文稿ID {prs_id} 未找到"

    # 验证文本颜色格式
    if text_color:
        if not text_color.startswith('#') or len(text_color) != 7:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"

        try:
            # 解析RGB值
            r = int(text_color[1:3], 16)
            g = int(text_color[3:5], 16)
            b = int(text_color[5:7], 16)
        except ValueError:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"

    # 验证对齐方式
    alignment_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY
    }

    if alignment.lower() not in alignment_map:
        return f"错误: 无效的对齐方式 '{alignment}'，支持的对齐方式有: {', '.join(alignment_map.keys())}"

    try:
        # 打开PowerPoint演示文稿
        prs = presentations[prs_id]

        # 获取幻灯片总数
        total_slides = len(prs.slides)

        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"

        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]

        # 添加文本框
        textbox = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )

        # 获取文本框的文本帧
        text_frame = textbox.text_frame

        # 设置文本内容
        text_frame.text = text

        # 设置文本格式
        p = text_frame.paragraphs[0]

        # 设置对齐方式
        p.alignment = alignment_map[alignment.lower()]

        # 获取run对象来设置字体属性
        run = p.runs[0]

        # 设置字体名称
        if font_name:
            run.font.name = font_name

        # 设置字体大小
        if font_size:
            run.font.size = Pt(font_size)

        # 设置字体加粗
        if font_bold:
            run.font.bold = True

        # 设置字体斜体
        if font_italic:
            run.font.italic = True

        # 设置文本颜色
        if text_color:
            run.font.color.rgb = RGBColor(r, g, b)

        return f"已成功在幻灯片 {slide_index} 添加文本框"

    except Exception as e:
        return f"添加文本框时出错: {str(e)}"


@mcp.tool()
def add_textboxes_bulk(prs_id: str, slide_index: int, boxes: List[Dict[str, Any]]) -> str:
    """
    一次性添加多个文本框。
    boxes: 每个字典表示一个文本框的参数（text、位置、尺寸、字体等）
    """
    messages = []
    for box in boxes:
        msg = add_text_box(prs_id, slide_index, **box)
        messages.append(msg)
    return "\n".join(messages)


# @mcp.tool()
# def insert_image(prs_id: str, slide_index: int, image_path: str, left: float = 1.0, top: float = 1.0,
#                  width: float = None, height: float = None) -> str:
#     """
#     在PowerPoint演示文稿的指定幻灯片中插入图片。
#
#     Args:
#         prs_id: PowerPoint演示文稿在内存中的唯一标识符
#         slide_index: 要插入图片的幻灯片索引（从1开始计数）
#         image_path: 图片文件的完整路径或相对路径
#         left: 图片左侧位置（英寸）
#         top: 图片顶部位置（英寸）
#         width: 图片宽度（英寸），如果不指定则保持原始宽高比
#         height: 图片高度（英寸），如果不指定则保持原始宽高比
#
#     Returns:
#         操作结果信息
#     """
#     # 检查是否安装了必要的库
#     if not pptx_installed:
#         return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
#
#     if prs_id not in presentations:
#         return f"错误: 演示文稿ID {prs_id} 未找到"
#
#     try:
#         # 打开PowerPoint演示文稿
#         prs = Presentation(prs_id)
#
#         # 获取幻灯片总数
#         total_slides = len(prs.slides)
#
#         # 检查slide_index是否有效
#         if slide_index < 1 or slide_index > total_slides:
#             return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
#
#         # 获取指定幻灯片
#         slide = prs.slides[slide_index - 1]
#
#         if image_path.startswith(('http://', 'https://')):
#             import requests
#             import tempfile
#
#             response = requests.get(image_path, stream=True)
#             response.raise_for_status()
#
#             img_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
#             img_temp_path = img_temp.name
#
#             with open(img_temp_path, 'wb') as f:
#                 for chunk in response.iter_content(1024):
#                     f.write(chunk)
#
#         else:
#             if not os.path.exists(image_path):
#                 return f"Error: Image file '{image_path}' not found"
#             else:
#                 img_temp_path = image_path
#
#
#         # 插入图片
#         if width and height:
#             # 按指定尺寸插入图片
#             slide.shapes.add_picture(
#                 img_temp_path,
#                 Inches(left),
#                 Inches(top),
#                 width=Inches(width),
#                 height=Inches(height)
#             )
#         elif width:
#             # 只指定宽度，高度按比例计算
#             slide.shapes.add_picture(
#                 img_temp_path,
#                 Inches(left),
#                 Inches(top),
#                 width=Inches(width)
#             )
#         elif height:
#             # 只指定高度，宽度按比例计算
#             slide.shapes.add_picture(
#                 img_temp_path,
#                 Inches(left),
#                 Inches(top),
#                 height=Inches(height)
#             )
#         else:
#             # 使用原始尺寸
#             slide.shapes.add_picture(
#                 img_temp_path,
#                 Inches(left),
#                 Inches(top)
#             )
#
#         return f"已成功在幻灯片 {slide_index} 插入图片"
#
#     except Exception as e:
#         return f"插入图片时出错: {str(e)}"
#
#
# @mcp.tool()
# def insert_table(prs_id: str, slide_index: int, rows: int, cols: int, data: List[List[str]] = None,
#                  left: float = 1.0, top: float = 1.0, width: float = 6.0, height: float = 3.0) -> str:
#     """
#     在PowerPoint演示文稿的指定幻灯片中插入表格。
#
#     Args:
#         prs_id: PowerPoint演示文稿在内存中的唯一标识符
#         slide_index: 要插入表格的幻灯片索引（从1开始计数）
#         rows: 表格行数
#         cols: 表格列数
#         data: 表格数据，二维列表，每个内部列表表示一行数据
#         left: 表格左侧位置（英寸）
#         top: 表格顶部位置（英寸）
#         width: 表格宽度（英寸）
#         height: 表格高度（英寸）
#
#     Returns:
#         操作结果信息
#     """
#     # 检查是否安装了必要的库
#     if not pptx_installed:
#         return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
#
#     if prs_id not in presentations:
#         return f"错误: 演示文稿ID {prs_id} 未找到"
#
#     # 验证行列数
#     if rows < 1 or cols < 1:
#         return "错误: 表格行数和列数必须大于0"
#
#     # 验证数据格式
#     if data:
#         if len(data) > rows:
#             return f"错误: 提供的数据行数({len(data)})超过了表格行数({rows})"
#
#         for i, row_data in enumerate(data):
#             if len(row_data) > cols:
#                 return f"错误: 第{i + 1}行的数据列数({len(row_data)})超过了表格列数({cols})"
#
#     try:
#         # 打开PowerPoint演示文稿
#         prs = Presentation(prs_id)
#
#         # 获取幻灯片总数
#         total_slides = len(prs.slides)
#
#         # 检查slide_index是否有效
#         if slide_index < 1 or slide_index > total_slides:
#             return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
#
#         # 获取指定幻灯片
#         slide = prs.slides[slide_index - 1]
#
#         # 创建表格形状
#         shape = slide.shapes.add_table(
#             rows,
#             cols,
#             Inches(left),
#             Inches(top),
#             Inches(width),
#             Inches(height)
#         )
#
#         # 获取表格对象
#         table = shape.table
#
#         # 填充表格数据
#         if data:
#             for i, row_data in enumerate(data):
#                 if i < rows:  # 确保不超过表格行数
#                     for j, cell_data in enumerate(row_data):
#                         if j < cols:  # 确保不超过表格列数
#                             # 获取单元格
#                             cell = table.cell(i, j)
#
#                             # 设置单元格文本
#                             cell.text = str(cell_data)
#
#         return f"已成功在幻灯片 {slide_index} 插入 {rows}×{cols} 的表格"
#
#     except Exception as e:
#         return f"插入表格时出错: {str(e)}"


@mcp.tool()
def add_slide_notes(prs_id: str, slide_index: int, notes_text: str) -> str:
    """
    为PowerPoint演示文稿的指定幻灯片添加备注。
    
    Args:
        prs_id: PowerPoint演示文稿在内存中的唯一标识符
        slide_index: 要添加备注的幻灯片索引（从1开始计数）
        notes_text: 备注内容
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    if prs_id not in presentations:
        return f"错误: 演示文稿ID {prs_id} 未找到"

    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(prs_id)

        # 获取幻灯片总数
        total_slides = len(prs.slides)

        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"

        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]

        # 获取备注页对象
        notes_slide = slide.notes_slide

        # 设置备注文本
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

        return f"已成功为幻灯片 {slide_index} 添加备注"

    except Exception as e:
        return f"添加幻灯片备注时出错: {str(e)}"


@mcp.tool()
def insert_shape(prs_id: str, slide_index: int, shape_type: str, left: float = 1.0, top: float = 1.0,
                 width: float = 2.0, height: float = 2.0, fill_color: str = None) -> str:
    """
    在PowerPoint演示文稿的指定幻灯片中插入形状。
    
    Args:
        prs_id: PowerPoint演示文稿在内存中的唯一标识符
        slide_index: 要插入形状的幻灯片索引（从1开始计数）
        shape_type: 形状类型，可选值包括: "rectangle", "oval", "triangle", "arrow", "line"
        left: 形状左侧位置（英寸）
        top: 形状顶部位置（英寸）
        width: 形状宽度（英寸）
        height: 形状高度（英寸）
        fill_color: 填充颜色的十六进制代码，例如"#FF0000"表示红色
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    if prs_id not in presentations:
        return f"错误: 演示文稿ID {prs_id} 未找到"

    # 验证形状类型
    shape_type_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.TRIANGLE,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "line": MSO_SHAPE.LINE
    }

    if shape_type.lower() not in shape_type_map:
        return f"错误: 无效的形状类型 '{shape_type}'，支持的形状类型有: {', '.join(shape_type_map.keys())}"

    # 验证颜色格式
    if fill_color:
        if not fill_color.startswith('#') or len(fill_color) != 7:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"

        try:
            # 解析RGB值
            r = int(fill_color[1:3], 16)
            g = int(fill_color[3:5], 16)
            b = int(fill_color[5:7], 16)
        except ValueError:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"

    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(prs_id)

        # 获取幻灯片总数
        total_slides = len(prs.slides)

        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"

        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]

        # 获取形状类型
        shape_type_value = shape_type_map[shape_type.lower()]

        # 添加形状
        shape = slide.shapes.add_shape(
            shape_type_value,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )

        # 设置填充颜色
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)


        return f"已成功在幻灯片 {slide_index} 插入 {shape_type} 形状"

    except Exception as e:
        return f"插入形状时出错: {str(e)}"


@mcp.tool()
def insert_chart(prs_id: str, slide_index: int, chart_type: str, data: List[List[str]], left: float = 1.0,
                 top: float = 1.0, width: float = 6.0, height: float = 4.0, has_legend: bool = True,
                 has_title: bool = True, title: str = "图表标题") -> str:
    """
    在PowerPoint演示文稿的指定幻灯片中插入图表。
    
    Args:
        prs_id: PowerPoint演示文稿在内存中的唯一标识
        slide_index: 要插入图表的幻灯片索引（从1开始计数）
        chart_type: 图表类型，可选值: "column", "line", "pie", "bar", "area"
        data: 图表数据，二维列表，第一行为类别标签，第一列为数据系列名称
        left: 图表左侧位置（英寸）
        top: 图表顶部位置（英寸）
        width: 图表宽度（英寸）
        height: 图表高度（英寸）
        has_legend: 是否显示图例
        has_title: 是否显示标题
        title: 图表标题文本
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    if prs_id not in presentations:
        return f"错误: 演示文稿ID {prs_id} 未找到"

    from pptx.enum.chart import XL_CHART_TYPE
    chart_type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "area": XL_CHART_TYPE.AREA
    }

    if chart_type.lower() not in chart_type_map:
        return f"错误: 无效的图表类型 '{chart_type}'，支持的图表类型有: {', '.join(chart_type_map.keys())}"

    # 验证数据格式
    if not data or len(data) < 2:
        return "错误: 图表数据必须至少包含两行（标题行和至少一个数据行）"

    for row in data:
        if not row or len(row) < 2:
            return "错误: 每行数据必须至少包含两列（类别名称和至少一个数据点）"

    try:
        # 打开PowerPoint演示文稿
        prs = presentations[prs_id]

        # 获取幻灯片总数
        total_slides = len(prs.slides)

        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"

        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]

        # 转换字符串数据为数值型数据
        # 第一行和第一列保持为字符串（类别和系列名称）
        # 其余部分转换为浮点数
        chart_data = []
        for i, row in enumerate(data):
            new_row = []
            for j, cell in enumerate(row):
                if i == 0 or j == 0:
                    # 第一行或第一列保持为字符串
                    new_row.append(str(cell))
                else:
                    # 尝试转换为浮点数
                    try:
                        new_row.append(float(cell))
                    except (ValueError, TypeError):
                        # 如果无法转换，使用0
                        new_row.append(0.0)
            chart_data.append(new_row)

        # 获取图表类型
        chart_type_value = chart_type_map[chart_type.lower()]

        # 创建图表
        chart_data_obj = None

        # 根据图表类型的不同，处理数据的方式也不同
        from pptx.chart.data import CategoryChartData, ChartData

        if chart_type.lower() == "pie":
            # 饼图只使用第二行数据
            chart_data_obj = ChartData()
            categories = [cat for cat in chart_data[0][1:]]
            values = [val for val in chart_data[1][1:]]

            # 添加类别和值
            chart_data_obj.categories = categories
            chart_data_obj.add_series(chart_data[1][0], values)
        else:
            # 其他图表类型
            chart_data_obj = CategoryChartData()

            # 添加类别（第一行除第一个单元格外的所有单元格）
            categories = [cat for cat in chart_data[0][1:]]
            chart_data_obj.categories = categories

            # 添加每个系列的数据（从第二行开始）
            for row in chart_data[1:]:
                series_name = row[0]
                values = [val for val in row[1:]]
                chart_data_obj.add_series(series_name, values)

        # 添加图表
        chart = slide.shapes.add_chart(
            chart_type_value,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            chart_data_obj
        ).chart

        # 设置图表标题
        if has_title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False

        # 设置图例
        chart.has_legend = has_legend
        if has_legend:
            chart.legend.position = 1  # 底部
            chart.legend.include_in_layout = False

        return f"已成功在幻灯片 {slide_index} 插入 {chart_type} 图表"

    except Exception as e:
        return f"插入图表时出错: {str(e)}"


# @mcp.tool()
# def apply_presentation_theme(file_path: str, theme_name: str) -> str:
#     """
#     为PowerPoint演示文稿应用主题。
#
#     Args:
#         file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
#         theme_name: 主题名称，可选值包括: "Office", "Adjacency", "Angles", "Apex",
#                    "Apothecary", "Aspect", "Austin", "Black Tie", "Civic", "Clarity",
#                    "Composite", "Concourse", "Couture", "Elemental", "Equity", "Essential",
#                    "Executive", "Grid", "Hardcover", "Horizon", "Integral", "Ion",
#                    "Ion Boardroom", "Median", "Metro", "Module", "Newsprint", "Opulent",
#                    "Organic", "Oriel", "Origin", "Paper", "Perspective", "Pushpin",
#                    "Slipstream", "Solstice", "Technic", "Thatch", "Trek", "Urban", "Vapor Trail", "Wisp"
#
#     Returns:
#         操作结果信息
#     """
#     # 检查是否安装了必要的库
#     if not pptx_installed:
#         return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
#
#     # 检查是否提供了完整路径
#     if not os.path.isabs(file_path):
#         # 从环境变量获取基础路径
#         base_path = os.environ.get('OFFICE_EDIT_PATH')
#         if not base_path:
#             base_path = os.path.join(os.path.expanduser('~'), '桌面')
#
#         # 构建完整路径
#         file_path = os.path.join(base_path, file_path)
#
#     # 确保文件存在
#     if not os.path.exists(file_path):
#         return f"错误: 文件 {file_path} 不存在"
#
#     # Python-pptx不直接支持更改主题，我们需要使用win32com
#     try:
#         import win32com.client
#
#         # 初始化PowerPoint应用程序
#         powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#         powerpoint.Visible = False  # 不显示PowerPoint窗口
#
#         # 打开演示文稿
#         presentation = powerpoint.Presentations.Open(file_path)
#
#         # 尝试应用主题
#         try:
#             # 获取主题文件的绝对路径
#             # 注意：这需要Office安装了对应的主题
#             # 主题文件通常位于Office的安装目录的THEMES子目录中
#             theme_file = f"{theme_name}.thmx"
#
#             # 应用主题
#             presentation.ApplyTheme(theme_file)
#
#             # 保存演示文稿
#             presentation.Save()
#             success = True
#             message = f"已成功应用主题 '{theme_name}'"
#         except Exception as theme_error:
#             success = False
#             message = f"应用主题时出错: {str(theme_error)}"
#
#         # 关闭演示文稿和PowerPoint应用程序
#         presentation.Close()
#         powerpoint.Quit()
#
#         if success:
#             return message
#         else:
#             # 如果使用COM方法失败，提供一个备选方法
#             return f"{message}\n可以尝试在PowerPoint中手动应用主题。"
#
#     except ImportError:
#         return "错误: 无法应用主题，这个功能需要在Windows系统上安装pywin32库。\n请使用命令安装: pip install pywin32"
#
#     except Exception as e:
#         return f"应用主题时出错: {str(e)}"
#
#
# @mcp.tool()
# def add_animation(file_path: str, slide_index: int, shape_index: int, animation_type: str,
#                   animation_trigger: str = "on_click", animation_delay: float = 0.0) -> str:
#     """
#     为PowerPoint演示文稿中的形状添加动画效果。
#
#     Args:
#         file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
#         slide_index: 要添加动画的幻灯片索引（从1开始计数）
#         shape_index: 要添加动画的形状索引（从1开始计数）
#         animation_type: 动画类型，可选值:
#                       "entrance": 入场动画
#                       "emphasis": 强调动画
#                       "exit": 退场动画
#         animation_trigger: 动画触发方式，可选值:
#                         "on_click": 点击时播放
#                         "with_previous": 与上一动画同时播放
#                         "after_previous": 在上一动画之后播放
#         animation_delay: 动画延迟时间（秒）
#
#     Returns:
#         操作结果信息
#     """
#     # 检查是否安装了必要的库
#     if not pptx_installed:
#         return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
#
#     # 检查是否提供了完整路径
#     if not os.path.isabs(file_path):
#         # 从环境变量获取基础路径
#         base_path = os.environ.get('OFFICE_EDIT_PATH')
#         if not base_path:
#             base_path = os.path.join(os.path.expanduser('~'), '桌面')
#
#         # 构建完整路径
#         file_path = os.path.join(base_path, file_path)
#
#     # 确保文件存在
#     if not os.path.exists(file_path):
#         return f"错误: 文件 {file_path} 不存在"
#
#     # 验证动画类型
#     if animation_type not in ["entrance", "emphasis", "exit"]:
#         return f"错误: 无效的动画类型 '{animation_type}'，支持的动画类型有: entrance, emphasis, exit"
#
#     # 验证动画触发方式
#     if animation_trigger not in ["on_click", "with_previous", "after_previous"]:
#         return f"错误: 无效的动画触发方式 '{animation_trigger}'，支持的触发方式有: on_click, with_previous, after_previous"
#
#     # Python-pptx不直接支持动画效果，我们需要使用win32com
#     try:
#         import win32com.client
#         from win32com.client import constants
#
#         # 初始化PowerPoint应用程序
#         powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#         powerpoint.Visible = False  # 不显示PowerPoint窗口
#
#         # 打开演示文稿
#         presentation = powerpoint.Presentations.Open(file_path)
#
#         # 检查幻灯片索引
#         if slide_index < 1 or slide_index > presentation.Slides.Count:
#             # 关闭文件和应用程序
#             presentation.Close()
#             powerpoint.Quit()
#             return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {presentation.Slides.Count} 张幻灯片"
#
#         # 获取指定幻灯片
#         slide = presentation.Slides.Item(slide_index)
#
#         # 检查形状索引
#         if shape_index < 1 or shape_index > slide.Shapes.Count:
#             # 关闭文件和应用程序
#             presentation.Close()
#             powerpoint.Quit()
#             return f"错误: 无效的形状索引 {shape_index}，幻灯片 {slide_index} 共有 {slide.Shapes.Count} 个形状"
#
#         # 获取指定形状
#         shape = slide.Shapes.Item(shape_index)
#
#         # 添加动画
#         animation_effect = None
#
#         # 获取动画序列
#         animation_sequence = slide.TimeLine.MainSequence
#
#         # 根据动画类型添加不同效果
#         if animation_type == "entrance":
#             # 入场动画
#             animation_effect = animation_sequence.AddEntrance(
#                 shape,
#                 constants.msoAnimEffectFade,  # 默认使用淡入效果
#                 constants.msoAnimateByObject,
#                 constants.msoAnimTriggerOnPageClick
#             )
#         elif animation_type == "emphasis":
#             # 强调动画
#             animation_effect = animation_sequence.AddEmphasis(
#                 shape,
#                 constants.msoAnimEffectPulse,  # 默认使用脉冲效果
#                 constants.msoAnimateByObject,
#                 constants.msoAnimTriggerOnPageClick
#             )
#         elif animation_type == "exit":
#             # 退场动画
#             animation_effect = animation_sequence.AddExit(
#                 shape,
#                 constants.msoAnimEffectFade,  # 默认使用淡出效果
#                 constants.msoAnimateByObject,
#                 constants.msoAnimTriggerOnPageClick
#             )
#
#         # 设置动画触发方式
#         if animation_trigger == "on_click":
#             animation_effect.Timing.TriggerType = constants.msoAnimTriggerOnPageClick
#         elif animation_trigger == "with_previous":
#             animation_effect.Timing.TriggerType = constants.msoAnimTriggerWithPrevious
#         elif animation_trigger == "after_previous":
#             animation_effect.Timing.TriggerType = constants.msoAnimTriggerAfterPrevious
#
#         # 设置动画延迟时间
#         animation_effect.Timing.TriggerDelayTime = animation_delay
#
#         # 保存演示文稿
#         presentation.Save()
#
#         # 关闭演示文稿和PowerPoint应用程序
#         presentation.Close()
#         powerpoint.Quit()
#
#         return f"已成功为幻灯片 {slide_index} 中的形状 {shape_index} 添加 {animation_type} 动画"
#
#     except ImportError:
#         return "错误: 无法添加动画，这个功能需要在Windows系统上安装pywin32库。\n请使用命令安装: pip install pywin32"
#
#     except Exception as e:
#         # 尝试关闭PowerPoint
#         try:
#             if 'presentation' in locals() and presentation:
#                 presentation.Close(False)
#             if 'powerpoint' in locals() and powerpoint:
#                 powerpoint.Quit()
#         except:
#             pass
#
#         return f"添加动画时出错: {str(e)}"
#
#
# @mcp.tool()
# def set_slide_transition(file_path: str, slide_index: int, transition_type: str, duration: float = 1.0,
#                          on_click: bool = True, auto_advance: bool = False, advance_time: float = 5.0) -> str:
#     """
#     设置PowerPoint演示文稿中幻灯片的切换效果。
#
#     Args:
#         file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
#         slide_index: 要设置切换效果的幻灯片索引（从1开始计数），0表示所有幻灯片
#         transition_type: 切换效果类型，可选值包括: "none", "fade", "push", "wipe", "split", "reveal",
#                        "random", "shape", "zoom", "gallery", "dissolve"
#         duration: 切换效果持续时间（秒）
#         on_click: 是否点击鼠标时切换幻灯片
#         auto_advance: 是否自动前进到下一张幻灯片
#         advance_time: 自动前进的时间（秒），仅当auto_advance=True时有效
#
#     Returns:
#         操作结果信息
#     """
#     # 检查是否安装了必要的库
#     if not pptx_installed:
#         return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
#
#     # 检查是否提供了完整路径
#     if not os.path.isabs(file_path):
#         # 从环境变量获取基础路径
#         base_path = os.environ.get('OFFICE_EDIT_PATH')
#         if not base_path:
#             base_path = os.path.join(os.path.expanduser('~'), '桌面')
#
#         # 构建完整路径
#         file_path = os.path.join(base_path, file_path)
#
#     # 确保文件存在
#     if not os.path.exists(file_path):
#         return f"错误: 文件 {file_path} 不存在"
#
#     # 验证切换效果类型
#     transition_types = ["none", "fade", "push", "wipe", "split", "reveal", "random", "shape", "zoom", "gallery",
#                         "dissolve"]
#     if transition_type.lower() not in transition_types:
#         return f"错误: 无效的切换效果类型 '{transition_type}'，支持的类型有: {', '.join(transition_types)}"
#
#     # Python-pptx不直接支持设置幻灯片切换效果，我们需要使用win32com
#     try:
#         import win32com.client
#         from win32com.client import constants
#
#         # 初始化PowerPoint应用程序
#         powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#         powerpoint.Visible = False  # 不显示PowerPoint窗口
#
#         # 打开演示文稿
#         presentation = powerpoint.Presentations.Open(file_path)
#
#         # 切换效果类型映射
#         transition_type_map = {
#             "none": constants.ppTransitionNone,
#             "fade": constants.ppTransitionFade,
#             "push": constants.ppTransitionPush,
#             "wipe": constants.ppTransitionWipe,
#             "split": constants.ppTransitionSplit,
#             "reveal": constants.ppTransitionReveal,
#             "random": constants.ppTransitionRandom,
#             "shape": constants.ppTransitionShapes,
#             "zoom": constants.ppTransitionZoom,
#             "gallery": constants.ppTransitionGallery,
#             "dissolve": constants.ppTransitionDissolve
#         }
#
#         # 获取切换效果类型的值
#         transition_value = transition_type_map.get(transition_type.lower())
#
#         # 确定要修改的幻灯片
#         if slide_index == 0:
#             # 修改所有幻灯片
#             for slide in presentation.Slides:
#                 # 设置切换效果
#                 slide.SlideShowTransition.EntryEffect = transition_value
#                 # 设置切换效果持续时间
#                 slide.SlideShowTransition.Duration = duration
#                 # 设置是否点击鼠标时切换
#                 slide.SlideShowTransition.AdvanceOnClick = on_click
#                 # 设置是否自动前进
#                 slide.SlideShowTransition.AdvanceOnTime = auto_advance
#                 # 设置自动前进的时间
#                 if auto_advance:
#                     slide.SlideShowTransition.AdvanceTime = advance_time
#
#             # 保存演示文稿
#             presentation.Save()
#
#             # 构建结果消息
#             result = f"已成功为所有幻灯片设置 {transition_type} 切换效果"
#         else:
#             # 修改指定幻灯片
#             # 检查幻灯片索引是否有效
#             if slide_index < 1 or slide_index > presentation.Slides.Count:
#                 # 关闭文件和应用程序
#                 presentation.Close()
#                 powerpoint.Quit()
#                 return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {presentation.Slides.Count} 张幻灯片"
#
#             # 获取指定幻灯片
#             slide = presentation.Slides.Item(slide_index)
#
#             # 设置切换效果
#             slide.SlideShowTransition.EntryEffect = transition_value
#             # 设置切换效果持续时间
#             slide.SlideShowTransition.Duration = duration
#             # 设置是否点击鼠标时切换
#             slide.SlideShowTransition.AdvanceOnClick = on_click
#             # 设置是否自动前进
#             slide.SlideShowTransition.AdvanceOnTime = auto_advance
#             # 设置自动前进的时间
#             if auto_advance:
#                 slide.SlideShowTransition.AdvanceTime = advance_time
#
#             # 保存演示文稿
#             presentation.Save()
#
#             # 构建结果消息
#             result = f"已成功为幻灯片 {slide_index} 设置 {transition_type} 切换效果"
#
#         # 关闭演示文稿和PowerPoint应用程序
#         presentation.Close()
#         powerpoint.Quit()
#
#         return result
#
#     except ImportError:
#         return "错误: 无法设置幻灯片切换效果，这个功能需要在Windows系统上安装pywin32库。\n请使用命令安装: pip install pywin32"
#
#     except Exception as e:
#         # 尝试关闭PowerPoint
#         try:
#             if 'presentation' in locals() and presentation:
#                 presentation.Close(False)
#             if 'powerpoint' in locals() and powerpoint:
#                 powerpoint.Quit()
#         except:
#             pass
#
#         return f"设置幻灯片切换效果时出错: {str(e)}"



@mcp.tool()
def load_presentation(file_path: str) -> str:
    """
    从文件加载一个现有的PowerPoint演示文稿到内存并读取其基本信息。

    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径

    Returns:
        演示文稿的基本信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法打开PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"

    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')

        # 构建完整路径
        file_path = os.path.join(base_path, file_path)

    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"

    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        prs_id = _generate_presentation_id(os.path.basename(file_path))

        # 确保ID唯一
        while prs_id in presentations:
            prs_id = _generate_presentation_id(os.path.basename(file_path))

        presentations[prs_id] = prs

        return get_presentation_info(prs_id)
    except Exception as e:
        return f"打开PowerPoint演示文稿时出错: {str(e)}"


@mcp.tool()
def save_to_oss(prs_id: str, oss_object_name: str = None) -> str:
    """
    将PowerPoint文件保存到OSS存储

    Args:
        file_path: 该演示文稿在内存中的唯一标识符
        oss_object_name: OSS上的对象名称（可选）

    Returns:
        OSS下载链接
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"

    prs = presentations[prs_id]

    try:
        # 生成临时文件名（使用UUID避免冲突）
        file_name = f"{prs_id}_{uuid.uuid4().hex[:8]}.pptx"
        tmp_path = os.path.join("/tmp", file_name)

        # 保存到临时文件
        prs.save(tmp_path)
        file_size = os.path.getsize(tmp_path)

        # 获取OSS配置（使用环境变量）
        endpoint = os.getenv("OSS_ENDPOINT")
        bucket_name = os.getenv("OSS_BUCKET_NAME")
        access_key = os.getenv("OSS_ACCESS_KEY")
        secret_key = os.getenv("OSS_SECRET_KEY")

        if not all([endpoint, bucket_name, access_key, secret_key]):
            return "OSS configuration missing"

        # 创建OSS客户端
        auth = oss2.Auth(access_key, secret_key)
        bucket = oss2.Bucket(auth, endpoint, bucket_name)

        # 上传文件
        oss_object_name = f"presentations/{file_name}"
        bucket.put_object_from_file(oss_object_name, tmp_path)

        # 生成可下载URL（1小时有效）
        download_url = bucket.sign_url('GET', oss_object_name, 3600)

        # 清理临时文件
        os.remove(tmp_path)

        return f"PPT已保存！下载地址：{download_url}"

    except Exception as e:
        error_msg = f"Error saving presentation: {str(e)}"
        print(error_msg, file=sys.stderr)
        return error_msg


def _generate_presentation_id(title: str) -> str:
    """生成唯一的演示文稿id"""
    cleaned_title = re.sub(r'[^\w\s\-]', '', title)
    cleaned_title = re.sub(r'[\s\-]+', '_', cleaned_title)
    cleaned_title = cleaned_title.strip('_').lower()
    if len(cleaned_title) > 20:
        cleaned_title = cleaned_title[:20].rstrip('_')
    short_uuid = uuid.uuid4().hex[:8]
    return f"{cleaned_title}_{short_uuid}" if cleaned_title else f"presentation_{short_uuid}"


def _format_text(frame_paragraph, text: str, set_color_to_white=False) -> None:
    r"""Apply bold and italic formatting while preserving the original
    word order.

    Args:
        frame_paragraph: The paragraph to format.
        text (str): The text to format.
        set_color_to_white (bool): Whether to set the color to white.
            (default: :obj:`False`)
    """

    from pptx.dml.color import RGBColor

    matches = list(BOLD_ITALICS_PATTERN.finditer(text))
    last_index = 0

    for match in matches:
        start, end = match.span()
        if start > last_index:
            run = frame_paragraph.add_run()
            run.text = text[last_index:start]
            if set_color_to_white:
                run.font.color.rgb = RGBColor(255, 255, 255)

        if match.group(2):  # Bold
            run = frame_paragraph.add_run()
            run.text = match.group(2)
            run.font.bold = True
            if set_color_to_white:
                run.font.color.rgb = RGBColor(255, 255, 255)
        elif match.group(3):  # Italics
            run = frame_paragraph.add_run()
            run.text = match.group(3)
            run.font.italic = True
            if set_color_to_white:
                run.font.color.rgb = RGBColor(255, 255, 255)

        last_index = end

    if last_index < len(text):
        run = frame_paragraph.add_run()
        run.text = text[last_index:]
        if set_color_to_white:
            run.font.color.rgb = RGBColor(255, 255, 255)


def _get_slide_width_height_inches(prs: "presentation.Presentation") -> Tuple[float, float]:

    r"""Get the dimensions of a slide in inches.

    Args:
        presentation (presentation.Presentation): The presentation object.

    Returns:
        Tuple[float, float]: The width and height in inches.
    """
    slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * (
        prs.slide_width or 0
    )
    slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * (
        prs.slide_height or 0
    )
    return slide_width_inch, slide_height_inch


def _remove_slide_number_from_heading(header: str) -> str:
    r"""Remove the slide number from a given slide header.

    Args:
        header (str): The header of a slide.

    Returns:
        str: The header without slide number.
    """
    if SLIDE_NUMBER_REGEX.match(header):
        idx = header.find(':')
        header = header[idx + 1:]
    return header

def _get_slide_placeholders(slide: "Slide") -> List[Tuple[int, str]]:
    r"""Return the index and name of all placeholders present in a slide.

    Args:
        slide (Slide): The slide.

    Returns:
        List[Tuple[int, str]]: A list containing placeholders (idx, name)
            tuples.
    """
    if hasattr(slide.shapes, 'placeholders'):
        placeholders = [
            (shape.placeholder_format.idx, shape.name.lower())
            for shape in slide.shapes.placeholders
        ]
        if placeholders and len(placeholders) > 0:
            placeholders.pop(0)  # Remove the title placeholder
        return placeholders
    return []

def _get_flat_list_of_contents(items: List[Union[str, List[Any]]], level: int) -> List[Tuple[str, int]]:
    r"""Flatten a hierarchical list of bullet points to a single list.

    Args:
        items (List[Union[str, List[Any]]]): A bullet point (string or
            list).
        level (int): The current level of hierarchy.

    Returns:
        List[Tuple[str, int]]: A list of (bullet item text, hierarchical
            level) tuples.
    """

    flat_list = []

    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list.extend(
                _get_flat_list_of_contents(item, level + 1)
            )

    return flat_list

def _add_bulleted_items(text_frame: "TextFrame", flat_items_list: List[Tuple[str, int]], set_color_to_white: bool = False) -> None:
    r"""Add a list of texts as bullet points and apply formatting.

    Args:
        text_frame (TextFrame): The text frame where text is to be
            displayed.
        flat_items_list (List[Tuple[str, int]]): The list of items to be
            displayed.
        set_color_to_white (bool): Whether to set the font color to white.
            (default: :obj:`False`)
    """
    if not flat_items_list:
        logger.warning("Empty bullet point list provided")
        return
    for idx, item_content in enumerate(flat_items_list):
        item_text, item_level = item_content

        if idx == 0:
            if not text_frame.paragraphs:
                # Ensure a paragraph exists if the frame is empty or
                # cleared
                paragraph = text_frame.add_paragraph()
            else:
                # Use the first existing paragraph
                paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()

        paragraph.level = item_level

        _format_text(
            paragraph,
            item_text.removeprefix(STEP_BY_STEP_PROCESS_MARKER),
            set_color_to_white=set_color_to_white,
        )



def main():
    try:
        mcp.run()
    except KeyboardInterrupt:
        print("服务器已停止")
    except Exception as e:
        print(f"服务器运行时出错: {str(e)}")
        sys.exit(1)





if __name__ == "__main__":
    main()
